"""주소 기반 자동 pptx 리포트 생성기.

report_template/generate_report.js(pptxgenjs 목업)를 실제 라이브 데이터로 채우는
python-pptx 버전이다. 두 부분으로 나뉜다:

  1. 데이터 계층: fetch_report_data() — building_example.py의 기존 함수들
     (build_master_report, REB 공실률, 소상공인 상가업소, 서울 상권분석)을 호출해
     주소 하나에 대한 리포트 데이터를 한 딕셔너리로 모은다.
  2. 생성 계층: generate_pptx(data) — 그 딕셔너리를 python-pptx로 그려 pptx
     바이트를 반환한다.

정성적 항목(개발호재/SWOT)은 어떤 공공API로도 자동 수집되지 않는다. 예전에는
report_template/manual_slides.pptx(연남동 클라이언트 리포트에서 그대로 복사한
슬라이드)를 모든 주소에 동일하게 끼워 넣었지만, 이러면 임의 주소로 리포트를
생성해도 연남동/홍대입구역 관련 가짜 문구가 그대로 남는 문제가 있었다. 지금은
대시보드 사이드바에 사용자가 직접 입력한 텍스트가 있을 때만 해당 슬라이드를
그리고, 비어 있으면 생략한다(가짜 데이터로 채우지 않음).
"""

import datetime
import io
import json
import math
import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.oxml.ns import qn

from building_example import (
    analyze_district_price_stats,
    analyze_seismic_risk,
    analyze_seoul_trade_area_detail,
    build_executive_summary,
    build_master_report,
    combine_zoning_sources,
    extract_year,
    get_building_ledger,
    _CORE_FIELD_LABELS,
    _MULTIROW_KEEP_COLS,
    _ONE_LINE_FIELDS,
    _clean,
    get_nearby_stores,
    get_reb_vacancy_snapshot,
    get_reb_vacancy_trend,
    get_seoul_trade_area_locations,
    get_seoul_trade_area_quarter_dataset,
    find_nearest_seoul_trade_area,
    reb_current_quarter_id,
    REB_COMMERCIAL_VACANCY_STATBL_IDS,
    SEOUL_TRDAR_SALES_SERVICE,
    SEOUL_TRDAR_STORE_SERVICE,
    SEOUL_TRDAR_FLPOP_SERVICE,
    SEOUL_TRDAR_WRC_POPLTN_SERVICE,
)
from ledger_ocr import extract_ledger_content

# ------------------------------------------------------------------
# 색상 팔레트 / 글꼴 — 딥 네이비(지배색) + 웜 테라코타 그라데이션 포인트.
# 부동산·시장데이터 리포트에 맞춰 무난한 파랑 대신 "무디 네이비 + 웜 테라코타"
# 조합을 골랐다(2025년 트렌드: 차분한 다크 베이스에 따뜻한 단색 포인트 하나만
# 강하게 주는 방식). 네이비가 화면의 대부분을 차지하는 지배색이고, 테라코타는
# 통계 숫자·차트 강조색으로만 아껴 쓰는 단일 포인트다.
# ------------------------------------------------------------------
NAVY = RGBColor(0x16, 0x21, 0x3E)
NAVY_DARK = RGBColor(0x0B, 0x12, 0x20)
NAVY_LIGHT = RGBColor(0x2A, 0x39, 0x5C)
ICE = RGBColor(0xE8, 0xEC, 0xF2)
TERRACOTTA = RGBColor(0xD9, 0x71, 0x3F)  # 단색 포인트(통계 숫자·강조 텍스트·차트 막대) — 실제 테라코타/코퍼 톤.
GRADIENT_A = RGBColor(0x2B, 0x4C, 0x7E)  # 그라데이션 시작 — 딥 네이비블루
GRADIENT_B = RGBColor(0xE2, 0x82, 0x5A)  # 그라데이션 끝 — 웜 테라코타/코럴
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x22, 0x26, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xF5, 0xF6, 0xF8)
GRID_LINE = RGBColor(0xDC, 0xE1, 0xE8)


def _gradient_fill(shape_or_chart_format, angle=45):
    """네이비블루→테라코타 그라데이션 채우기. 배지·포인트바·차트 시리즈 등 강조 도형에 쓴다."""
    fill = shape_or_chart_format.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = GRADIENT_A
    stops[0].position = 0.0
    stops[-1].color.rgb = GRADIENT_B
    stops[-1].position = 1.0
    fill.gradient_angle = angle

# Malgun Gothic을 bold로 쓰면 이 프로젝트가 QA에 쓰는 LibreOffice 렌더러에서 한글이
# 깨진 장식체로 나오는 버그가 있었다(report_template 작업 중 발견) — 실제 PowerPoint에서도
# 재현될 수 있어 위험하므로, 정상체/굵게 모두 깨끗하게 렌더링되는 Calibri로 통일한다.
# 한글은 OS 기본 대체 글꼴(Windows에서는 사실상 맑은 고딕)로 표시된다.
FONT_NAME = "Calibri"

# 모든 슬라이드 함수의 x/y/w/h 리터럴은 원래 16:9(13.333×7.5in) 캔버스 기준으로
# 짜여 있다. 실제 출력 크기를 A4 가로(297×210mm)로 바꾸면서 각 슬라이드 함수의
# 좌표 자체는 그대로 두고, Inches()로 변환되는 시점에만 가로/세로 각각 다른
# 배율(SCALE_X/SCALE_Y)을 곱해 새 캔버스에 여백 없이 꽉 차도록 늘린다 — 배율을
# 하나로 통일하면(가로 기준으로만 맞추면) A4가 16:9보다 더 정방형에 가까운
# 비율이라 아래쪽에 빈 여백이 크게 남는다. 원(뱃지·장식 동그라미)처럼 가로세로
# 비율이 유지돼야 하는 도형만 SCALE_U(더 작은 쪽 배율)로 균일하게 스케일한다.
DESIGN_W_IN = 13.333
DESIGN_H_IN = 7.5
SLIDE_W_IN = 11.69  # A4 가로 297mm
SLIDE_H_IN = 8.27   # A4 가로 210mm
SCALE_X = SLIDE_W_IN / DESIGN_W_IN
SCALE_Y = SLIDE_H_IN / DESIGN_H_IN
SCALE_U = min(SCALE_X, SCALE_Y)


def _sx(v):
    """디자인 그리드의 x좌표/가로폭 값을 실제 인치로 변환."""
    return Inches(v * SCALE_X)


def _sy(v):
    """디자인 그리드의 y좌표/세로높이 값을 실제 인치로 변환."""
    return Inches(v * SCALE_Y)


def _su(v):
    """원·정사각 배지처럼 가로세로 비율을 유지해야 하는 도형의 크기 변환."""
    return Inches(v * SCALE_U)


# ==================================================================
# 생성 계층 (python-pptx)
# ==================================================================
def _new_slide(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY if dark else WHITE
    return slide


def _textbox(slide, x, y, w, h, text, size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False, font=FONT_NAME,
             line_spacing=None):
    box = slide.shapes.add_textbox(_sx(x), _sy(y), _sx(w), _sy(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def _rich_textbox(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    """runs: [(text, {size,bold,color}), ...] 한 줄에 여러 스타일을 섞어 쓸 때."""
    box = slide.shapes.add_textbox(_sx(x), _sy(y), _sx(w), _sy(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for text, opts in runs:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 12))
        run.font.bold = opts.get("bold", False)
        run.font.name = FONT_NAME
        run.font.color.rgb = opts.get("color", TEXT_DARK)
    return box


def _section_title(slide, text):
    _textbox(slide, 0.6, 0.45, 11.3, 0.6, text, size=28, bold=True, color=NAVY)


def _page_footer(slide, address, label):
    _textbox(slide, 0.6, 7.12, 10, 0.3, f"{address}  ·  {label}", size=9, color=MUTED)


def _card(slide, x, y, w, h, fill=CARD_BG, line=GRID_LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(x), _sy(y), _sx(w), _sy(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _add_picture_placeholder(slide, x, y, w, h, image_stream=None, idx=90):
    """PowerPoint에서 드래그&드롭으로 바로 교체할 수 있는 실제 사진 placeholder를 만든다.
    image_stream이 있으면(자동 생성 지도 등) 미리 채워 넣고, 없으면 빈 상자로 둔다 — 어느 쪽이든
    <p:ph type="pic">를 갖는 진짜 placeholder라서 파일탐색기에서 이미지를 끌어다 놓으면 그대로 바뀐다."""
    if image_stream is not None:
        shape = slide.shapes.add_picture(image_stream, _sx(x), _sy(y), _sx(w), _sy(h))
        nvPr = shape._element.find(qn("p:nvPicPr")).find(qn("p:nvPr"))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _sx(x), _sy(y), _sx(w), _sy(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = GRID_LINE
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        nvPr = shape._element.find(qn("p:nvSpPr")).find(qn("p:nvPr"))
    nvPr.append(nvPr.makeelement(qn("p:ph"), {"type": "pic", "idx": str(idx)}))
    return shape


def _stat_card(slide, x, y, w, h, value, label, sub=None):
    _card(slide, x, y, w, h)
    _textbox(slide, x + 0.15, y + 0.12, w - 0.3, h * 0.45, value, size=22, bold=True,
             color=TERRACOTTA, valign=MSO_ANCHOR.BOTTOM)
    _textbox(slide, x + 0.15, y + h * 0.55, w - 0.3, h * 0.22, label, size=11, bold=True, color=TEXT_DARK)
    if sub:
        _textbox(slide, x + 0.15, y + h * 0.76, w - 0.3, h * 0.2, sub, size=9, color=MUTED)


def _key_point_box(slide, text, x=0.6, y=6.15, w=12.1, h=0.75):
    """참고 리포트(연남동 상권분석)의 시그니처 'KEY POINT' 콜아웃 — 그 슬라이드에서 데이터가
    말해주는 한 줄 통찰을 짙은 네이비 박스 + 테라코타 라벨로 눈에 띄게 강조한다. 각 분석
    슬라이드가 이미 계산해 둔 note/insight 문장을 그대로 여기 태워서 쓴다(새 계산 없음)."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(x), _sy(y), _sx(w), _sy(h))
    box.adjustments[0] = 0.14
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.shadow.inherit = False
    box.text_frame.paragraphs[0].text = ""

    badge_w = 1.35
    _textbox(slide, x + 0.3, y, badge_w, h, "KEY POINT", size=10, bold=True,
             color=TERRACOTTA, valign=MSO_ANCHOR.MIDDLE)
    _textbox(slide, x + 0.3 + badge_w, y + 0.08, w - badge_w - 0.6, h - 0.16, text,
             size=12, color=WHITE, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.2)


def _table(slide, x, y, w, h, headers, rows, col_ratios=None, font_size=11):
    n_rows, n_cols = len(rows) + 1, len(headers)
    gframe = slide.shapes.add_table(n_rows, n_cols, _sx(x), _sy(y), _sx(w), _sy(h))
    table = gframe.table
    if col_ratios:
        total = sum(col_ratios)
        for i, ratio in enumerate(col_ratios):
            table.columns[i].width = _sx(w * ratio / total)

    def _fmt_cell(cell, text, bold, color, fill):
        cell.text = "" if text is None else str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_top = cell.margin_bottom = Pt(4)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = FONT_NAME
            p.font.color.rgb = color

    for j, htext in enumerate(headers):
        _fmt_cell(table.cell(0, j), htext, True, WHITE, NAVY)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            _fmt_cell(table.cell(i, j), val, False, TEXT_DARK, WHITE)
    return table


def _style_chart_axes(chart):
    chart.has_legend = False
    chart.has_title = False
    try:
        chart.category_axis.tick_labels.font.size = Pt(9.5)
        chart.category_axis.tick_labels.font.color.rgb = MUTED
        chart.category_axis.format.line.color.rgb = GRID_LINE
        chart.value_axis.tick_labels.font.size = Pt(9.5)
        chart.value_axis.tick_labels.font.color.rgb = MUTED
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = GRID_LINE
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass


def _bar_chart(slide, x, y, w, h, categories, values, color, horizontal=False, num_fmt="0", label_size=9):
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    chart_data.add_series("값", values)
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gframe = slide.shapes.add_chart(chart_type, _sx(x), _sy(y), _sx(w), _sy(h), chart_data)
    chart = gframe.chart
    _style_chart_axes(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.font.size = Pt(label_size)
    dl.font.color.rgb = TEXT_DARK
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()
    return chart


def _line_chart(slide, x, y, w, h, categories, values, color, num_fmt="0.0", label_size=9):
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    chart_data.add_series("값", values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, _sx(x), _sy(y), _sx(w), _sy(h), chart_data)
    chart = gframe.chart
    _style_chart_axes(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.font.size = Pt(label_size)
    dl.font.color.rgb = NAVY
    series = plot.series[0]
    series.format.line.color.rgb = color
    series.format.line.width = Pt(2.5)
    series.marker.style = 8  # circle
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = color
    series.smooth = False
    return chart


def _cover_slide(prs, data):
    slide = _new_slide(prs, dark=True)
    for cx, cy, cw, ch, color, grad in [
        (8.6, 3.0, 7.5, 7.5, NAVY_LIGHT, False),
        (9.6, 4.0, 5.5, 5.5, NAVY_DARK, False),
        (10.4, 4.8, 3.9, 3.9, None, True),
    ]:
        ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(cx), _sy(cy), _su(cw), _su(ch))
        if grad:
            _gradient_fill(ellipse, angle=45)
        else:
            ellipse.fill.solid()
            ellipse.fill.fore_color.rgb = color
        ellipse.line.fill.background()
        ellipse.shadow.inherit = False

    _textbox(slide, 0.9, 2.35, 8, 0.4, "부동산 분석 리포트", size=14, bold=True, color=TERRACOTTA)
    _textbox(slide, 0.9, 2.8, 9.5, 1.6, data["address"], size=40, bold=True, color=WHITE)
    _textbox(slide, 0.9, 4.05, 8, 0.5, "건축물 · 실거래가 · 상권 통합 분석", size=16, color=ICE)

    stats = data.get("cover_stats") or []
    for i, s in enumerate(stats[:3]):
        cx = 0.9 + i * 2.9
        _textbox(slide, cx, 5.5, 2.6, 0.55, s["value"], size=24, bold=True, color=TERRACOTTA)
        _textbox(slide, cx, 6.05, 2.6, 0.5, s["label"], size=10.5, color=ICE)

    _textbox(slide, 0.9, 6.7, 6, 0.4, data["report_date"], size=11, color=MUTED)


def _toc_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "목차")
    _textbox(slide, 0.6, 1.05, 11, 0.3, f"{data['address']} 분석 리포트 구성", size=12, color=MUTED)

    parts = data["toc_parts"]
    col_w, gap_x, start_x, start_y = 3.95, 0.2, 0.6, 1.55
    page_no = 1
    for pi, part in enumerate(parts):
        x = start_x + pi * (col_w + gap_x)
        _textbox(slide, x, start_y, col_w, 0.5, part["title"], size=13.5, bold=True, color=TERRACOTTA)
        iy = start_y + 0.55
        for item in part["items"]:
            badge_w = 0.4 if page_no >= 10 else 0.32
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(x), _sy(iy), _su(badge_w), _su(0.32))
            badge.adjustments[0] = 0.2
            badge.fill.solid()
            badge.fill.fore_color.rgb = NAVY
            badge.line.fill.background()
            badge.shadow.inherit = False
            tf = badge.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(page_no)
            run.font.size = Pt(9 if page_no >= 10 else 10)
            run.font.bold = True
            run.font.name = FONT_NAME
            run.font.color.rgb = WHITE
            _textbox(slide, x + 0.5, iy - 0.03, col_w - 0.53, 0.38, item, size=11.5, valign=MSO_ANCHOR.MIDDLE)
            iy += 0.44
            page_no += 1
    _page_footer(slide, data["address"], "목차")


def _section_divider_slide(prs, part_no, kicker, title, subtitle):
    """PART 구분용 다크 배경 슬라이드. 예전에는 manual_slides.pptx의 연남동 전용 슬라이드를
    그대로 복사했지만, 지금은 실제 법정동명을 넣은 일반 문장만 쓰는 코드 생성 버전이다."""
    slide = _new_slide(prs, dark=True)
    ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(9.5), _sy(-2.0), _su(6.5), _su(6.5))
    ellipse.fill.solid()
    ellipse.fill.fore_color.rgb = NAVY_LIGHT
    ellipse.line.fill.background()
    ellipse.shadow.inherit = False

    _textbox(slide, 0.9, 2.2, 6, 1.6, part_no, size=64, bold=True, color=NAVY_LIGHT)
    _textbox(slide, 0.9, 3.55, 8, 0.4, kicker, size=14, bold=True, color=TERRACOTTA)
    _textbox(slide, 0.9, 3.95, 10, 1.0, title, size=34, bold=True, color=WHITE)
    _textbox(slide, 0.9, 4.95, 9.5, 0.8, subtitle, size=15, color=ICE, line_spacing=1.3)


def _summary_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "핵심 요약")
    _card(slide, 0.6, 1.25, 12.1, 1.35)
    _textbox(slide, 0.9, 1.42, 11.5, 1.0, data["summary_text"], size=13, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.25)

    stats = data["summary_stats"]
    cols, gap = 3, 0.25
    card_w = (12.1 - gap * (cols - 1)) / cols
    card_h = 1.35
    start_y = 2.85
    for i, s in enumerate(stats):
        col, row = i % cols, i // cols
        x = 0.6 + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        _stat_card(slide, x, y, card_w, card_h, s["value"], s["label"], s.get("sub"))
    _page_footer(slide, data["address"], "핵심 요약")


def _building_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "건축물 개요")
    b = data["building"]

    left_x, left_y, left_w, left_h = 0.6, 1.3, 6.7, 5.6
    _card(slide, left_x, left_y, left_w, left_h)
    core = b["core"]
    row_h = (left_h - 0.4) / max(1, math.ceil(len(core) / 2))
    for i, (label, value) in enumerate(core):
        col, row = i % 2, i // 2
        x = left_x + 0.3 + col * (left_w / 2 - 0.15)
        y = left_y + 0.25 + row * row_h
        _textbox(slide, x, y, left_w / 2 - 0.5, row_h * 0.42, label, size=10.5, color=MUTED)
        _textbox(slide, x, y + row_h * 0.38, left_w / 2 - 0.5, row_h * 0.5, value, size=15, bold=True, color=NAVY)

    right_x, right_w = 7.5, 5.2
    next_y = 1.3
    if b.get("floors"):
        _textbox(slide, right_x, next_y, right_w, 0.35, "층별 면적 (㎡)", size=13, bold=True, color=NAVY)
        chart = _bar_chart(
            slide, right_x, next_y + 0.4, right_w, 2.7,
            [f["label"] for f in b["floors"]], [f["value"] for f in b["floors"]],
            TERRACOTTA, horizontal=True,
        )
        _gradient_fill(chart.plots[0].series[0].format, angle=0)
        next_y += 3.35

    if b.get("zoning"):
        _textbox(slide, right_x, next_y, right_w, 0.35, "용도지역 · 지구", size=13, bold=True, color=NAVY)
        py = next_y + 0.4
        for z in b["zoning"][:5]:
            pw = min(right_w, 0.35 + len(z) * 0.16)
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(right_x), _sy(py), _sx(pw), _sy(0.42))
            pill.adjustments[0] = 0.5
            pill.fill.solid()
            pill.fill.fore_color.rgb = ICE
            pill.line.fill.background()
            pill.shadow.inherit = False
            tf = pill.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = z
            run.font.size = Pt(10.5)
            run.font.name = FONT_NAME
            run.font.color.rgb = NAVY
            py += 0.55
    _page_footer(slide, data["address"], "건축물 개요")


def _seismic_age_slide(prs, data):
    seismic = data.get("seismic")
    if not seismic or not seismic.get("dong_dist"):
        return
    slide = _new_slide(prs)
    _section_title(slide, "노후도 · 내진 분석")

    _textbox(slide, 0.6, 1.3, 6.6, 0.35, "동 전체 내진설계 분류 분포", size=13, bold=True, color=NAVY)
    chart = _bar_chart(
        slide, 0.6, 1.7, 6.6, 3.6,
        [d["label"] for d in seismic["dong_dist"]], [d["value"] for d in seismic["dong_dist"]],
        NAVY_LIGHT, horizontal=True,
    )

    if seismic.get("subject_label"):
        _stat_card(slide, 7.5, 1.3, 5.2, 1.5, seismic["subject_label"], "대상 건물 내진 분류", seismic.get("subject_full"))

    if seismic.get("age_dist"):
        _textbox(slide, 7.5, 3.1, 5.2, 0.35, "노후도 구간별 분포 (동 전체)", size=13, bold=True, color=NAVY)
        _bar_chart(slide, 7.5, 3.5, 5.2, 2.9,
                   [d["label"] for d in seismic["age_dist"]], [d["value"] for d in seismic["age_dist"]], TERRACOTTA)
    _page_footer(slide, data["address"], "노후도 · 내진 분석")


def _intro_divider_slide(prs, data):
    dong = data.get("location", {}).get("adong_name") or ""
    subtitle = (
        f"{dong} 일대의 입지 여건과 주변 상권 구조를 실제 데이터로 살펴본다" if dong
        else "입지 여건과 주변 상권 구조를 실제 데이터로 살펴본다"
    )
    _section_divider_slide(prs, "01", "PART 1", "입지 및 상권 분석", subtitle)


def _location_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "위치 및 입지")
    loc = data["location"]

    map_x, map_y, map_w, map_h = 0.6, 1.3, 7.6, 5.6
    _add_picture_placeholder(slide, map_x, map_y, map_w, map_h, image_stream=loc.get("map_image"), idx=90)
    if not loc.get("map_image"):
        _textbox(slide, map_x, map_y + map_h / 2 - 0.35, map_w, 0.7,
                 "지도 이미지를 생성하지 못했습니다.\nPowerPoint에서 이 영역에 이미지를 끌어다 놓으면 바로 채워집니다.",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    rx, rw = 8.5, 4.2
    _card(slide, rx, 1.3, rw, 1.3)
    _textbox(slide, rx + 0.25, 1.42, rw - 0.5, 0.3, "행정동", size=10.5, color=MUTED)
    _textbox(slide, rx + 0.25, 1.68, rw - 0.5, 0.4, loc.get("adong_name") or "-", size=18, bold=True, color=NAVY)
    _textbox(slide, rx + 0.25, 2.1, rw - 0.5, 0.3, "법정동", size=10.5, color=MUTED)
    _textbox(slide, rx + 0.25, 2.36, rw - 0.5, 0.2, loc.get("ldong_name") or "-", size=13)

    _textbox(slide, rx, 2.85, rw, 0.35, "인근 지하철역", size=13, bold=True, color=NAVY)
    sy = 3.25
    for st_ in loc.get("subway", [])[:3]:
        _card(slide, rx, sy, rw, 0.85, fill=CARD_BG)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(rx + 0.18), _sy(sy + 0.2), _su(0.44), _su(0.44))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor.from_string(st_.get("color", "1E2A44").lstrip("#"))
        badge.line.color.rgb = WHITE
        badge.line.width = Pt(1.5)
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = st_["badge"]
        run.font.size = Pt(8.5 if len(st_["badge"]) > 1 else 12)
        run.font.bold = True
        run.font.name = FONT_NAME
        run.font.color.rgb = WHITE
        _textbox(slide, rx + 0.78, sy + 0.1, rw - 1.0, 0.35, st_["name"], size=11.5, bold=True)
        _textbox(slide, rx + 0.78, sy + 0.45, rw - 1.0, 0.3, f"{st_['line']} · 도보 약 {st_['dist']}", size=9.5, color=MUTED)
        sy += 1.0
    _page_footer(slide, data["address"], "위치 및 입지")


def _transactions_slide(prs, data):
    tx = data.get("transactions")
    if not tx or not tx.get("rows"):
        return
    slide = _new_slide(prs)
    _section_title(slide, "실거래가 동향")
    _textbox(slide, 0.6, 1.3, 5.6, 0.35, "최근 거래 내역 (반경 200m 이내)", size=13, bold=True, color=NAVY)
    _table(slide, 0.6, 1.7, 5.6, min(2.6, 0.42 * (len(tx["rows"]) + 1)),
           ["거래일", "층", "전용면적", "거래가격", "거리"], tx["rows"], col_ratios=[1.4, 0.9, 1.2, 1.3, 0.9])

    if tx.get("trend") and len(tx["trend"]) >= 2:
        _textbox(slide, 6.6, 1.3, 6.1, 0.35, "거래가 추이", size=13, bold=True, color=NAVY)
        _line_chart(slide, 6.6, 1.7, 6.1, 2.7, [t["label"] for t in tx["trend"]], [t["value"] for t in tx["trend"]], TERRACOTTA)

    if tx.get("insight") or tx.get("note"):
        _key_point_box(slide, tx.get("insight") or tx["note"])
    _page_footer(slide, data["address"], "실거래가 동향")


def _land_vs_floor_price_slide(prs, data):
    """참고 리포트(연남동 상권분석) '대지기준 vs 연면적기준 평당가 비교' 슬라이드와 같은 방식.

    반경 200m 실거래 매칭 결과 중 대지면적·연면적이 모두 확인되는 거래만 골라, 같은 거래를
    두 단위로 환산해 토지가치·건물가치 관점을 분리해 보여준다. 새 API 호출 없이 실거래가
    동향 슬라이드와 같은 데이터를 재계산만 한다(_build_transactions의 land_vs_floor)."""
    rows = [r for r in (data.get("transactions") or {}).get("land_vs_floor") or []
            if r["land_price"] is not None and r["floor_price"] is not None]
    if not rows:
        return
    slide = _new_slide(prs)
    _section_title(slide, "대지기준 vs 연면적기준 평당가 비교")
    _textbox(slide, 0.6, 1.05, 11.5, 0.3,
             "같은 거래를 두 가지 단위로 환산 — 토지가치와 건물가치의 분리 관점 (반경 200m 이내 매칭 거래)",
             size=11, color=MUTED)

    def _price_label(v):
        return f"{v / 1e4:,.2f}억" if v >= 1e4 else f"{v:,.0f}만원"

    table_rows = [
        [r["지번"], f"{r['far']}%" if r["far"] is not None else "-",
         _price_label(r["land_price"]), _price_label(r["floor_price"])]
        for r in rows
    ]
    _table(slide, 0.6, 1.55, 12.1, min(4.6, 0.45 * (len(table_rows) + 1)),
           ["지번", "용적률", "토지평당가", "건물평당가"], table_rows, col_ratios=[2.4, 1.3, 1.8, 1.8])

    scored = [(r, r["floor_price"] / r["land_price"]) for r in rows if r["land_price"]]
    if scored:
        top_r, ratio = max(scored, key=lambda t: abs(t[1] - 1))
        far_label = f"용적률 {top_r['far']}%" if top_r["far"] is not None else "용적률 정보 없음"
        if ratio > 1:
            insight = (f"{top_r['지번']}는 연면적기준가가 대지기준가의 {ratio:.1f}배로, "
                       f"저층·저용적({far_label}) 필지일수록 토지가치 중심 거래로 보입니다.")
        else:
            insight = (f"{top_r['지번']}는 연면적기준가가 대지기준가의 {ratio * 100:.0f}%에 불과해, "
                       f"고용적률({far_label}) 필지는 넓은 연면적이 단위가격을 희석시키는 경향을 보입니다.")
        _key_point_box(slide, insight, y=6.25, h=0.7)
    _page_footer(slide, data["address"], "대지 vs 연면적 평당가 비교")


def _district_slide(prs, data):
    dist = data.get("district")
    if not dist or not dist.get("age_buckets"):
        return
    slide = _new_slide(prs)
    _section_title(slide, f"동단위 시장 통계 — {data['location'].get('adong_name') or ''}")
    _textbox(slide, 0.6, 1.3, 8, 0.35, "준공연도대별 건물 수 분포 (동 전체)", size=13, bold=True, color=NAVY)
    _bar_chart(slide, 0.6, 1.7, 8.0, 4.3,
               [b["label"] for b in dist["age_buckets"]], [b["value"] for b in dist["age_buckets"]], NAVY_LIGHT)

    rx, rw = 9.0, 3.7
    if dist.get("callout"):
        _stat_card(slide, rx, 1.7, rw, 2.2, dist["callout"]["value"], dist["callout"]["label"], dist["callout"].get("sub"))
    if dist.get("note"):
        _key_point_box(slide, dist["note"])
    _page_footer(slide, data["address"], "동단위 시장 통계")


def _district_mix_slide(prs, data):
    dist = data.get("district")
    if not dist or not dist.get("mix"):
        return
    slide = _new_slide(prs)
    _section_title(slide, f"동단위 시장 통계 — 주용도별 구성 ({data['location'].get('adong_name') or ''})")
    _textbox(slide, 0.6, 1.3, 6.0, 0.35, "주용도별 건물 수 비중", size=13, bold=True, color=NAVY)
    chart = _bar_chart(slide, 0.6, 1.7, 6.0, 4.6,
                        [m["label"] for m in dist["mix"]], [m["value"] for m in dist["mix"]],
                        NAVY_LIGHT, horizontal=True, num_fmt="0.0")

    if dist.get("benchmark"):
        _textbox(slide, 6.9, 1.3, 5.8, 0.35, "주용도별 규모 벤치마크 (중앙값)", size=13, bold=True, color=NAVY)
        _table(slide, 6.9, 1.7, 5.8, min(2.6, 0.42 * (len(dist["benchmark"]) + 1)),
               ["주용도", "층수", "용적률(%)", "높이(m)"], dist["benchmark"], col_ratios=[2.2, 1, 1.3, 1.3])
    _page_footer(slide, data["address"], "동단위 시장 통계 — 주용도별")


def _price_history_slide(prs, data):
    ph = data.get("price_history")
    if not ph or not ph.get("trend"):
        return
    slide = _new_slide(prs)
    _section_title(slide, "공시가격 시계열")
    _textbox(slide, 0.6, 1.3, 7.4, 0.35, "연도별 공시가격 추이", size=13, bold=True, color=NAVY)
    _line_chart(slide, 0.6, 1.7, 7.4, 4.2, [t["label"] for t in ph["trend"]], [t["value"] for t in ph["trend"]], TERRACOTTA)

    if ph.get("rows"):
        _textbox(slide, 8.3, 1.3, 4.4, 0.35, "연도별 변동률", size=13, bold=True, color=NAVY)
        _table(slide, 8.3, 1.7, 4.4, min(2.0, 0.42 * (len(ph["rows"]) + 1)),
               ["연도", "공시가격", "전년대비"], ph["rows"], col_ratios=[1.3, 1.6, 1.5])
    district = ph.get("district")
    if district:
        _textbox(slide, 8.3, 3.9, 4.4, 0.3, "동단위 주변사례 비교", size=12, bold=True, color=NAVY)
        cards = [
            ("평균", f"{district['평균']:.2f}억" if district.get("평균") is not None else "-"),
            ("중앙값", f"{district['중앙값']:.2f}억" if district.get("중앙값") is not None else "-"),
        ]
        if district.get("상위백분율") is not None:
            cards.append(("대상물건 위치", f"상위 {district['상위백분율']:.0f}%"))
        cw, gap = (4.4 - 0.15 * (len(cards) - 1)) / len(cards), 0.15
        for i, (label, val) in enumerate(cards):
            _stat_card(slide, 8.3 + i * (cw + gap), 4.25, cw, 1.6, val, label)
    if ph.get("note"):
        _key_point_box(slide, ph["note"], y=6.05, h=0.85)
    _page_footer(slide, data["address"], "공시가격 시계열")


def _commercial_slide(prs, data):
    com = data.get("commercial")
    if not com or (not com.get("vacancy_trend") and not com.get("top_industries")):
        return
    slide = _new_slide(prs)
    area_name = com.get("area_name")
    title = f"상권 개황 — 공실률 및 주변 업종 ({area_name})" if area_name else "상권 개황 — 공실률 및 주변 업종"
    _section_title(slide, title)

    if com.get("vacancy_trend"):
        _textbox(slide, 0.6, 1.3, 6.0, 0.35, f"공실률 추이 ({com.get('vacancy_label', '')})", size=13, bold=True, color=NAVY)
        _line_chart(slide, 0.6, 1.7, 6.0, 4.2,
                    [t["label"] for t in com["vacancy_trend"]], [t["value"] for t in com["vacancy_trend"]], NAVY_LIGHT)
    else:
        _card(slide, 0.6, 1.7, 6.0, 4.2)
        _textbox(slide, 0.6, 3.5, 6.0, 0.6, "해당 지역 공실률 데이터를 찾지 못했습니다.",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER)

    if com.get("top_industries"):
        industry_label = "업종 Top 5 (점포 수, 서울 열린데이터광장)" if area_name else "반경 500m 업종 Top 5 (점포 수)"
        _textbox(slide, 6.9, 1.3, 5.8, 0.35, industry_label, size=13, bold=True, color=NAVY)
        chart = _bar_chart(slide, 6.9, 1.7, 5.8, 4.2,
                            [t["label"] for t in com["top_industries"]], [t["value"] for t in com["top_industries"]],
                            TERRACOTTA, horizontal=True)
        _gradient_fill(chart.plots[0].series[0].format, angle=0)
    if com.get("insight"):
        _key_point_box(slide, com["insight"], y=6.05, h=0.85)
    _page_footer(slide, data["address"], "상권 개황")


def _nearby_stores_detail_slide(prs, data):
    com = data.get("commercial") or {}
    stores = com.get("store_list")
    if not stores:
        return
    slide = _new_slide(prs)
    cols = com.get("store_list_cols") or ["상호명", "업종", "주소", "거리"]
    is_seoul_detail = cols[0] == "업종"
    area_name = com.get("area_name")
    if is_seoul_detail:
        title = f"업종별 상세 현황 — {area_name} (서울 열린데이터광장)" if area_name else "업종별 상세 현황"
        ratios = [2.6, 1.6, 1.3, 1.3, 1.3]
    else:
        title = "반경 상가업소 상세 목록"
        ratios = [2.2, 1.6, 3.4, 1]
    _section_title(slide, title)
    _table(slide, 0.6, 1.3, 12.1, min(5.6, 0.42 * (len(stores) + 1)), cols, stores, col_ratios=ratios)
    _page_footer(slide, data["address"], "업종별 상세 현황" if is_seoul_detail else "반경 상가업소 상세")


def _trade_area_map_slide(prs, data):
    tam = data.get("trade_area_map")
    if not tam:
        return
    slide = _new_slide(prs)
    _section_title(slide, f"상권영역 지도 — {tam.get('name', '')}")
    map_x, map_y, map_w, map_h = 0.6, 1.3, 12.1, 5.4
    _add_picture_placeholder(slide, map_x, map_y, map_w, map_h, image_stream=tam.get("map_image"), idx=91)
    _textbox(slide, map_x, map_y + map_h + 0.05, map_w, 0.3,
             "※ OpenStreetMap 기반 위치 참고 지도. 상권영역 경계 자체는 서울시 API가 중심좌표+면적만 제공해 "
             "정확한 폴리곤으로 표시할 수 없습니다. PowerPoint에서 이 영역에 이미지를 끌어다 놓으면 바로 교체됩니다.",
             size=8.5, italic=True, color=MUTED)
    _page_footer(slide, data["address"], "상권영역 지도")


def _seoul_detail_slide(prs, data):
    sd = data.get("seoul_trade_area")
    if not sd:
        return
    slide = _new_slide(prs)
    _section_title(slide, f"서울 상권 상세 — {sd.get('name', '')}")

    stats = sd.get("stats") or []
    cols, gap = 3, 0.25
    card_w = (12.1 - gap * (cols - 1)) / cols
    for i, s in enumerate(stats[:3]):
        x = 0.6 + i * (card_w + gap)
        _stat_card(slide, x, 1.3, card_w, 1.2, s["value"], s["label"])

    if sd.get("top_industries"):
        _textbox(slide, 0.6, 2.85, 5.6, 0.35, "업종별 매출 · 점포수 Top 5", size=13, bold=True, color=NAVY)
        _table(slide, 0.6, 3.25, 5.6, 2.85, ["업종", "당월 매출", "점포수"], sd["top_industries"], col_ratios=[3.0, 2.0, 1.6])

    if sd.get("weekday"):
        _textbox(slide, 6.6, 2.85, 6.1, 0.35, "요일별 매출 (억원)", size=13, bold=True, color=NAVY)
        chart = _bar_chart(slide, 6.6, 3.25, 6.1, 2.85, [d["label"] for d in sd["weekday"]], [d["value"] for d in sd["weekday"]],
                            TERRACOTTA, num_fmt="0.0")
        _gradient_fill(chart.plots[0].series[0].format, angle=0)
    if sd.get("insights"):
        _key_point_box(slide, sd["insights"][0], y=6.25, h=0.7)
    _page_footer(slide, data["address"], "서울 상권 상세")


def _seoul_detail_slide2(prs, data):
    """서울 상권 상세의 두 번째 슬라이드 — 성별·연령대·주중/주말 소비자 특성 +

    점포 개폐업 현황 + data.seoul.go.kr 데이터 기반 규칙형 인사이트. 첫 슬라이드가
    총괄·업종·요일별 매출을 다뤘다면, 이 슬라이드는 '누가, 언제, 어떤 강도로'
    소비하는 상권인지를 보여준다.
    """
    sd = data.get("seoul_trade_area")
    if not sd or not (sd.get("store") or sd.get("age") or sd.get("gender") or sd.get("insights")):
        return
    slide = _new_slide(prs)
    _section_title(slide, f"서울 상권 상세 — 소비자 특성 및 점포 ({sd.get('name', '')})")

    store = sd.get("store") or {}
    cards = [
        ("총점포수", store.get("총점포수"), "개"),
        ("프랜차이즈점포수", store.get("프랜차이즈점포수"), "개"),
        ("개업률", store.get("개업률(%)"), "%"),
        ("폐업률", store.get("폐업률(%)"), "%"),
    ]
    cards = [c for c in cards if c[1] is not None]
    if cards:
        cols, gap = len(cards), 0.2
        card_w = (12.1 - gap * (cols - 1)) / cols
        for i, (label, val, unit) in enumerate(cards):
            x = 0.6 + i * (card_w + gap)
            value_str = f"{val:,.0f}{unit}" if unit == "개" else f"{val}{unit}"
            _stat_card(slide, x, 1.3, card_w, 1.1, value_str, label)

    age = sd.get("age") or []
    if age:
        _textbox(slide, 0.6, 2.7, 5.6, 0.35, "연령대별 매출 비중(%)", size=13, bold=True, color=NAVY)
        chart = _bar_chart(slide, 0.6, 3.1, 5.6, 3.4, [a["label"] for a in age], [a["value"] for a in age],
                            TERRACOTTA, num_fmt="0.0")
        _gradient_fill(chart.plots[0].series[0].format, angle=0)

    gender = sd.get("gender") or {}
    weekend = sd.get("weekend_mix") or {}
    y = 2.7
    if gender:
        _textbox(slide, 6.6, y, 6.1, 0.3, "성별 매출 비중", size=13, bold=True, color=NAVY)
        _textbox(slide, 6.6, y + 0.35, 6.1, 0.5,
                 f"남성 {gender['남성비율(%)']:.0f}%   ·   여성 {gender['여성비율(%)']:.0f}%",
                 size=14, color=TEXT_DARK)
        y += 0.95
    if weekend:
        _textbox(slide, 6.6, y, 6.1, 0.3, "주중/주말 매출 비중", size=13, bold=True, color=NAVY)
        _textbox(slide, 6.6, y + 0.35, 6.1, 0.5,
                 f"주중 {weekend['주중비율(%)']:.0f}%   ·   주말 {weekend['주말비율(%)']:.0f}%",
                 size=14, color=TEXT_DARK)
        y += 0.95

    insights = sd.get("insights") or []
    if insights:
        _textbox(slide, 6.6, y + 0.1, 6.1, 0.3, "🔎 자동 인사이트", size=13, bold=True, color=NAVY)
        y += 0.5
        for line in insights[:4]:
            _textbox(slide, 6.6, y, 6.1, 0.5, f"• {line}", size=11, color=TEXT_DARK, line_spacing=1.2)
            y += 0.5
    _page_footer(slide, data["address"], "서울 상권 상세")


def _seoul_population_slide(prs, data):
    """서울 열린데이터광장 우리마을가게 상권분석서비스의 생활인구(유동인구)·직장인구를

    성별·연령대별로 시각화한다. 상권 매출/업종 데이터와 같은 API 호출(_build_seoul)에서
    이미 받아오는 값인데 지금까지 슬라이드에 안 쓰이고 있었다 — 새 API 호출 없이 기존
    데이터를 그대로 채워 넣는다."""
    sd = data.get("seoul_trade_area")
    if not sd or not (sd.get("flpop") or sd.get("wrc")):
        return
    slide = _new_slide(prs)
    _section_title(slide, f"유동인구(생활인구) · 직장인구 분석 — {sd.get('name', '')}")
    _textbox(slide, 0.6, 1.05, 11.5, 0.3,
             "※ 서울 열린데이터광장 우리마을가게 상권분석서비스 기준, 해당 상권 전체 분기 추정치입니다.",
             size=10, italic=True, color=MUTED)

    flpop = sd.get("flpop") or {}
    wrc = sd.get("wrc") or {}
    col_w, gap = 5.75, 0.6
    cols = [
        (0.6, "유동인구 (생활인구)", flpop, sd.get("flpop_age") or [], "총생활인구"),
        (0.6 + col_w + gap, "직장인구", wrc, sd.get("wrc_age") or [], "총직장인구"),
    ]
    for x, label, pop_stats, age_chart, total_key in cols:
        _textbox(slide, x, 1.55, col_w, 0.35, label, size=14, bold=True, color=NAVY)
        y = 2.0
        total = pop_stats.get(total_key)
        if total is not None:
            male = pop_stats.get("남성비율(%)")
            female = pop_stats.get("여성비율(%)")
            sub = f"남성 {male:.0f}% · 여성 {female:.0f}%" if male is not None and female is not None else None
            _stat_card(slide, x, y, col_w, 1.2, f"{total:,.0f}명", "분기 추정 인구", sub)
            y += 1.4
        if age_chart:
            _textbox(slide, x, y, col_w, 0.3, "연령대별 인구", size=12, bold=True, color=NAVY)
            chart = _bar_chart(slide, x, y + 0.35, col_w, 2.9,
                                [a["label"] for a in age_chart], [a["value"] for a in age_chart],
                                NAVY_LIGHT if "생활" in label else TERRACOTTA, num_fmt="#,##0")
            _gradient_fill(chart.plots[0].series[0].format, angle=0)
        elif total is None:
            _card(slide, x, y, col_w, 2.0)
            _textbox(slide, x, y + 0.85, col_w, 0.4, "데이터 없음", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    _page_footer(slide, data["address"], "유동인구 · 직장인구 분석")


def _growth_drivers_slide(prs, data):
    items = data.get("growth_drivers") or []
    if not items:
        return
    slide = _new_slide(prs)
    _section_title(slide, "개발호재 종합")
    _textbox(slide, 0.6, 1.05, 11.5, 0.3,
             "※ 사이드바에 직접 입력한 내용입니다 — 공공데이터로 자동 생성되지 않았습니다.",
             size=10, italic=True, color=MUTED)
    y = 1.55
    for i, item in enumerate(items[:8]):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(0.6), _sy(y + 0.05), _su(0.36), _su(0.36))
        _gradient_fill(badge, angle=45)
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.name = FONT_NAME
        run.font.color.rgb = WHITE
        _textbox(slide, 1.15, y, 11.35, 0.65, item, size=13, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
        y += 0.72
    _page_footer(slide, data["address"], "개발호재 종합")


_SWOT_META = [
    ("strengths", "S", "강점", GRADIENT_A),
    ("weaknesses", "W", "약점", TERRACOTTA),
    ("opportunities", "O", "기회", NAVY_LIGHT),
    ("threats", "T", "위협", NAVY_DARK),
]


def _swot_slide(prs, data):
    swot = data.get("swot") or {}
    if not any(swot.get(key) for key, *_ in _SWOT_META):
        return
    slide = _new_slide(prs)
    _section_title(slide, "SWOT 분석")
    _textbox(slide, 0.6, 1.05, 11.5, 0.3,
             "※ 사이드바에 직접 입력한 내용입니다 — 공공데이터로 자동 생성되지 않았습니다.",
             size=10, italic=True, color=MUTED)
    cell_w, cell_h, gap = 5.95, 2.55, 0.2
    x0, y0 = 0.6, 1.5
    for idx, (key, letter, kr, color) in enumerate(_SWOT_META):
        items = swot.get(key) or []
        if not items:
            continue
        col, row = idx % 2, idx // 2
        x, y = x0 + col * (cell_w + gap), y0 + row * (cell_h + gap)
        _card(slide, x, y, cell_w, cell_h)
        _textbox(slide, x + 0.2, y + 0.15, cell_w - 0.4, 0.4, f"{letter} — {kr}", size=14, bold=True, color=color)
        body = "\n".join(f"· {t}" for t in items[:5])
        _textbox(slide, x + 0.2, y + 0.6, cell_w - 0.4, cell_h - 0.8, body, size=11, line_spacing=1.25)
    _page_footer(slide, data["address"], "SWOT 분석")


def _key_metrics_table_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "핵심지표 종합 비교표")

    core_map = dict(data.get("building", {}).get("core") or [])
    rows = [
        ["대지면적", core_map.get("대지면적", "데이터 없음")],
        ["연면적", core_map.get("연면적", "데이터 없음")],
        ["사용승인일", core_map.get("사용승인일", "데이터 없음")],
    ]
    seismic = data.get("seismic") or {}
    rows.append(["대상 건물 내진 분류", seismic.get("subject_label") or "데이터 없음"])

    tx_rows = data.get("transactions", {}).get("rows") or []
    rows.append(["최근 실거래가", tx_rows[0][3] if tx_rows else "데이터 없음"])

    ph_rows = data.get("price_history", {}).get("rows") or []
    rows.append(["최근 공시가격", ph_rows[0][1] if ph_rows else "데이터 없음"])

    com = data.get("commercial") or {}
    rows.append(["공실률", f"{com['vacancy_trend'][-1]['value']}%" if com.get("vacancy_trend") else "데이터 없음"])

    seoul = data.get("seoul_trade_area")
    rows.append(["상권 추정매출(월)", seoul["stats"][0]["value"] if seoul and seoul.get("stats") else "데이터 없음"])

    _table(slide, 0.6, 1.3, 12.1, min(5.6, 0.5 * (len(rows) + 1)), ["항목", "값"], rows, col_ratios=[1, 2], font_size=13)
    _page_footer(slide, data["address"], "핵심지표 종합 비교표")


_DATA_SOURCES = [
    "국토교통부 건축HUB 건축물대장정보 서비스 (표제부 · 주택가격 · 지역지구구역 등)",
    "국토교통부 실거래가 공개시스템 (PublicDataReader 경유)",
    "한국부동산원 R-ONE 부동산통계정보시스템 (중대형 상가 공실률)",
    "서울 열린데이터광장 우리마을가게 상권분석서비스 (업종별 매출 · 점포수 · 소비자특성, 서울 소재 주소만 해당)",
    "소상공인시장진흥공단 상가(상권)정보 Open API (서울 열린데이터광장을 쓸 수 없을 때의 대체 데이터)",
    "카카오맵 Local API (주소 지오코딩)",
    "업로드한 건축물대장 열람본 이미지 OCR (소유자현황 · 변동사항 · 위반건축물 상세)",
]

_APPENDIX_DISCLAIMER = (
    "· 본 리포트는 공공데이터 API 응답을 자동 집계한 참고 자료이며, 법적 효력이 있는 감정평가 · 중개 문서가 아닙니다.\n"
    "· 공시가격은 시세가 아닙니다(통상 시세의 60~70% 수준, 연도별 현실화율 정책 변동 포함). 동단위 주변사례 평균 · "
    "중앙값 · 백분위는 참고용 비교치일 뿐 감정평가액이 아닙니다.\n"
    "· 개발호재 · SWOT 항목은 작성자가 사이드바에 직접 입력한 내용이 있을 때만 표시되며, 자동 생성되지 않습니다.\n"
    "· 실거래가는 대상 물건 좌표 기준 반경 200m 이내 거래를 지오코딩으로 매칭한 결과입니다. 좌표 확보(지오코딩)에 "
    "실패한 거래는 반경 계산에서 제외될 수 있습니다.\n"
    "· 상권 데이터는 서울 열린데이터광장 상권코드 단위로 집계되며, 대상 주소가 속한 상권 전체의 통계입니다 — "
    "개별 건물 단위 수치가 아닙니다."
)


def _address_ledger_content_slide(prs, data):
    """업로드된 건축물대장 열람본 이미지를 그대로 붙이는 대신, OCR로 추출한
    소유자현황·변동사항·관계자·위반건축물 내용을 이 리포트의 다른 슬라이드와 같은
    표/텍스트 형식으로 정리해서 보여준다. generate_ledger_pptx()의 열람본 전용
    리포트가 이미 쓰고 있는 _ledger_owner_history_slide()를 그대로 재사용한다."""
    content = data.get("ledger_content")
    if not content:
        return
    _ledger_owner_history_slide(
        prs, data["address"], content["owners"], content["changes"],
        firms=content.get("firms"), builder=content.get("builder"),
        violation_detail=content.get("violation_detail"),
    )


def _appendix_slide(prs, data):
    slide = _new_slide(prs)
    _section_title(slide, "부록 — 데이터 출처 및 유의사항")
    _textbox(slide, 0.6, 1.3, 11.5, 0.35, "데이터 출처", size=13, bold=True, color=NAVY)
    _textbox(slide, 0.6, 1.7, 11.5, 2.1, "\n".join(f"· {s}" for s in _DATA_SOURCES), size=12, line_spacing=1.35)
    _textbox(slide, 0.6, 4.0, 11.5, 0.35, "유의사항", size=13, bold=True, color=NAVY)
    _textbox(slide, 0.6, 4.4, 11.5, 2.4, _APPENDIX_DISCLAIMER, size=11, color=MUTED, line_spacing=1.4)
    _page_footer(slide, data["address"], "부록")


def _conclusion_slide(prs, data):
    slide = _new_slide(prs, dark=True)
    ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(-2.5), _sy(4.2), _su(6.5), _su(6.5))
    ellipse.fill.solid()
    ellipse.fill.fore_color.rgb = NAVY_LIGHT
    ellipse.line.fill.background()
    ellipse.shadow.inherit = False

    _textbox(slide, 0.9, 0.7, 8, 0.7, "종합 의견", size=32, bold=True, color=WHITE)
    cy = 1.9
    for i, point in enumerate(data.get("conclusion") or []):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(0.9), _sy(cy + 0.05), _su(0.36), _su(0.36))
        _gradient_fill(badge, angle=45)
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.name = FONT_NAME
        run.font.color.rgb = WHITE
        _textbox(slide, 1.5, cy, 10.6, 0.9, point, size=15, color=ICE, line_spacing=1.3)
        cy += 1.15
    _textbox(slide, 0.9, 6.9, 10, 0.4, f"{data['address']}  ·  {data['report_date']}", size=10, color=MUTED)


def _slide_plan(data):
    """실제로 그려질 슬라이드만 순서대로 나열한 (파트번호, 목차라벨, 렌더함수) 리스트.

    표지/목차 자신은 여기 포함하지 않는다(별도 처리). 각 항목의 존재 여부(조건부
    슬라이드)를 여기 한 곳에서만 판정해서, 목차 페이지 번호가 실제 슬라이드 순서와
    어긋나는 일이 구조적으로 생기지 않게 한다 — 예전에는 목차 항목 리스트를
    fetch_report_data()에서 슬라이드 스킵 조건과 별개로 하드코딩해서, 데이터가 없어
    슬라이드가 조용히 생략될 때마다 그 뒤 모든 페이지 번호가 밀리는 버그가 있었다."""
    dist = data.get("district") or {}
    com = data.get("commercial") or {}
    swot = data.get("swot") or {}

    return [
        (1, "핵심 요약", _summary_slide),
        (1, "건축물 개요", _building_slide),
        (1, "노후도 · 내진 분석", _seismic_age_slide) if data.get("seismic", {}).get("dong_dist") else None,
        (1, "입지 · 상권 개관", _intro_divider_slide),
        (1, "위치 및 입지", _location_slide),
        (2, "실거래가 동향", _transactions_slide) if data.get("transactions", {}).get("rows") else None,
        (2, "대지기준 vs 연면적기준 평당가 비교", _land_vs_floor_price_slide)
        if any(r["land_price"] is not None and r["floor_price"] is not None
               for r in data.get("transactions", {}).get("land_vs_floor") or [])
        else None,
        (2, "동단위 시장 통계 — 연대별", _district_slide) if dist.get("age_buckets") else None,
        (2, "동단위 시장 통계 — 주용도별", _district_mix_slide) if dist.get("mix") else None,
        (2, "공시가격 시계열", _price_history_slide) if data.get("price_history", {}).get("trend") else None,
        (2, "상권 개황", _commercial_slide) if (com.get("vacancy_trend") or com.get("top_industries")) else None,
        (2, "상권 업종 상세", _nearby_stores_detail_slide) if com.get("store_list") else None,
        (2, "상권영역 지도", _trade_area_map_slide) if data.get("trade_area_map") else None,
        (2, "서울 상권 상세", _seoul_detail_slide) if data.get("seoul_trade_area") else None,
        (2, "서울 상권 상세 — 소비자 특성", _seoul_detail_slide2)
        if (data.get("seoul_trade_area") or {}).get("insights") or (data.get("seoul_trade_area") or {}).get("age")
        else None,
        (2, "유동인구 · 직장인구 분석", _seoul_population_slide)
        if (data.get("seoul_trade_area") or {}).get("flpop") or (data.get("seoul_trade_area") or {}).get("wrc")
        else None,
        (3, "개발호재 종합", _growth_drivers_slide) if data.get("growth_drivers") else None,
        (3, "SWOT 분석", _swot_slide) if any(swot.get(k) for k in ("strengths", "weaknesses", "opportunities", "threats")) else None,
        (3, "핵심지표 종합 비교표", _key_metrics_table_slide),
        (3, "종합 의견", _conclusion_slide),
        (3, "소유자현황 · 변동사항 (OCR)", _address_ledger_content_slide) if data.get("ledger_content") else None,
        (3, "부록 — 데이터 출처 및 유의사항", _appendix_slide),
    ]


def generate_pptx(data: dict) -> bytes:
    """fetch_report_data()가 만든 딕셔너리로 pptx를 그려 바이트로 반환."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    plan = [item for item in _slide_plan(data) if item is not None]

    part_titles = {1: "PART 01 · 입지 및 건축물 분석", 2: "PART 02 · 시장 데이터 분석", 3: "PART 03 · 종합 평가"}
    data["toc_parts"] = [
        {"title": part_titles[part_no], "items": ["표지", "목차"] + [label for p, label, _ in plan if p == part_no]}
        if part_no == 1 else
        {"title": part_titles[part_no], "items": [label for p, label, _ in plan if p == part_no]}
        for part_no in (1, 2, 3)
    ]

    _cover_slide(prs, data)
    _toc_slide(prs, data)
    for _part_no, _label, render_fn in plan:
        render_fn(prs, data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
# ==================================================================
# 데이터 계층 — 주소 -> 실제 데이터
# ==================================================================
_SUBWAY_STATIONS_PATH = os.path.join(os.path.dirname(__file__), "address_map_component", "subway_stations.json")


def _load_subway_stations():
    with open(_SUBWAY_STATIONS_PATH, encoding="utf-8") as f:
        return json.load(f)


def _haversine_m(lat1, lon1, lat2, lon2):
    r = 6371000
    p1, p2 = math.radians(lat1), math.radians(lat2)
    dphi = math.radians(lat2 - lat1)
    dl = math.radians(lon2 - lon1)
    a = math.sin(dphi / 2) ** 2 + math.cos(p1) * math.cos(p2) * math.sin(dl / 2) ** 2
    return r * 2 * math.asin(math.sqrt(a))


def _nearby_subway(lat, lon, limit=3):
    stations = _load_subway_stations()
    scored = sorted(
        ((_haversine_m(lat, lon, s_lat, s_lon), name, lines) for name, s_lat, s_lon, lines in stations),
        key=lambda t: t[0],
    )
    out = []
    for d, name, lines in scored[:limit]:
        line_names = [l[0] for l in lines]
        color = lines[0][1].lstrip("#")
        badge = line_names[0][:2] if len(line_names[0]) > 2 else line_names[0]
        out.append({
            "name": f"{name}역", "line": "·".join(line_names), "color": color,
            "dist": f"{int(round(d / 10) * 10)}m", "badge": badge,
        })
    return out


def _render_location_map(lat, lon, subway=None, radius_m=None):
    """staticmap(OSM 타일, 브라우저 불필요)으로 위치 지도 PNG를 만들어 BytesIO로 반환."""
    try:
        from staticmap import StaticMap, CircleMarker
    except ImportError:
        return None
    try:
        m = StaticMap(1300, 1040, url_template="https://a.tile.openstreetmap.org/{z}/{x}/{y}.png")
        for st_ in (subway or [])[:3]:
            # 지하철 마커는 위치 참고용 점만 (역 좌표는 subway_stations.json 기준, 정밀 좌표 아님)
            pass
        m.add_marker(CircleMarker((lon, lat), "#FFFFFF", 30))
        m.add_marker(CircleMarker((lon, lat), "#C1652F", 22))
        m.add_marker(CircleMarker((lon, lat), "#FFFFFF", 8))
        image = m.render(zoom=16)
        buf = io.BytesIO()
        image.save(buf, format="PNG")
        buf.seek(0)
        return buf
    except Exception:
        return None


def _tx_field(row, candidates):
    for c in candidates:
        if c in row.index and pd.notna(row[c]) and str(row[c]).strip() not in ("", "nan"):
            return row[c]
    return None


def _format_tx_date(row):
    for c in ["계약일", "거래일"]:
        v = _tx_field(row, [c])
        if v:
            return str(v)
    y = _tx_field(row, ["년", "dealYear", "계약년도"])
    m = _tx_field(row, ["월", "dealMonth", "계약월"])
    d = _tx_field(row, ["일", "dealDay", "계약일자"])
    if y and m:
        try:
            return f"{int(y)}.{int(m):02d}" + (f".{int(d):02d}" if d else "")
        except (ValueError, TypeError):
            pass
    return "-"


_PYEONG_M2 = 3.305785  # 1평(平) = 약 3.305785㎡


def _land_vs_floor_price_rows(df, max_rows=8):
    """대지면적·연면적이 둘 다 있는 매칭 거래에 한해 토지평당가·건물평당가를 계산.

    같은 거래를 대지 기준/연면적 기준 두 단위로 환산해 토지가치·건물가치 관점을
    분리해서 보여준다 — 참고 리포트(연남동 상권분석)의 "대지기준 vs 연면적기준
    평당가 비교" 슬라이드와 같은 방식. 반경 매칭으로 이미 받아온 거래 데이터를
    그대로 재계산할 뿐, 추가 API 호출은 없다."""
    if df is None or df.empty:
        return []
    out = []
    for _, row in df.iterrows():
        jibun = _tx_field(row, ["지번"]) or "-"
        price = _tx_field(row, ["거래금액", "물건금액"])
        land = _tx_field(row, ["대지면적"])
        floor = _tx_field(row, ["연면적"])
        price_num = pd.to_numeric(str(price).replace(",", ""), errors="coerce") if price is not None else None
        land_num = pd.to_numeric(land, errors="coerce") if land is not None else None
        floor_num = pd.to_numeric(floor, errors="coerce") if floor is not None else None
        if pd.isna(price_num) or (pd.isna(land_num) and pd.isna(floor_num)):
            continue
        land_price = float(price_num) / (float(land_num) / _PYEONG_M2) if land_num and land_num > 0 else None
        floor_price = float(price_num) / (float(floor_num) / _PYEONG_M2) if floor_num and floor_num > 0 else None
        far = round(float(floor_num) / float(land_num) * 100) if land_num and floor_num and land_num > 0 else None
        out.append({
            "지번": str(jibun), "land_price": land_price, "floor_price": floor_price, "far": far,
        })
    out.sort(key=lambda r: r["land_price"] or 0, reverse=True)
    return out[:max_rows]


def _build_transactions(master, months_lookback):
    tx = master.get("실거래가") or {}
    result = {"rows": [], "trend": [], "note": tx.get("note"), "insight": None, "land_vs_floor": []}
    df = tx.get("df")
    if tx.get("status") != "matched" or df is None or df.empty:
        return result

    try:
        rows = []
        trend = []
        sort_col = None
        for c in ["계약일", "거래일", "년", "dealYear"]:
            if c in df.columns:
                sort_col = c
                break
        has_distance = "거리(m)" in df.columns
        sorted_df = df.sort_values(sort_col, ascending=False) if sort_col else df
        for _, row in sorted_df.head(8).iterrows():
            date_label = _format_tx_date(row)
            floor = _tx_field(row, ["층"])
            area = _tx_field(row, ["전용면적", "대지면적", "연면적", "계약면적"])
            price = _tx_field(row, ["거래금액", "물건금액"])
            price_num = pd.to_numeric(str(price).replace(",", ""), errors="coerce") if price is not None else None
            price_label = f"{price_num / 1e4:.1f}억" if pd.notna(price_num) else (str(price) if price else "-")
            dist = row.get("거리(m)") if has_distance else None
            dist_label = f"{int(dist)}m" if pd.notna(dist) else "-"
            rows.append([
                date_label,
                f"{floor}층" if floor not in (None, "") else "-",
                f"{float(area):.1f}㎡" if area not in (None, "") and pd.notna(pd.to_numeric(area, errors="coerce")) else "-",
                price_label,
                dist_label,
            ])
            if pd.notna(price_num):
                trend.append({"label": date_label, "value": round(float(price_num) / 1e4, 1)})
        result["rows"] = rows
        result["trend"] = list(reversed(trend))

        if len(result["trend"]) >= 2:
            first, last = result["trend"][0]["value"], result["trend"][-1]["value"]
            change_pct = (last / first - 1) * 100 if first else 0
            direction = "상승" if change_pct >= 5 else ("하락" if change_pct <= -5 else "보합")
            result["insight"] = (
                f"반경 200m 이내 거래가는 {result['trend'][0]['label']} {first:.1f}억에서 "
                f"{result['trend'][-1]['label']} {last:.1f}억으로 {direction} 흐름입니다 ({change_pct:+.1f}%)."
            )
        elif result["rows"]:
            result["insight"] = f"반경 200m 이내 최근 {months_lookback}개월간 {len(result['rows'])}건의 거래가 확인됩니다."

        result["land_vs_floor"] = _land_vs_floor_price_rows(df)
    except Exception:
        pass
    return result


def _build_district(master):
    stats = master.get("동단위통계")
    old_df = master.get("노후도")
    subject_age = int(old_df.iloc[0]["경과연수"]) if old_df is not None and not old_df.empty else None

    result = {"age_buckets": [], "callout": None, "note": None, "mix": [], "benchmark": []}
    if not stats or stats.get("연대별") is None or stats["연대별"].empty:
        return result

    by_decade = stats["연대별"]
    result["age_buckets"] = [
        {"label": row["연대"], "value": int(row["건수"])} for _, row in by_decade.iterrows()
    ]
    avg_age = (stats.get("총괄") or {}).get("평균경과연수")
    if avg_age is not None:
        if subject_age is not None:
            cmp_word = "높은" if subject_age > avg_age else "낮은"
            result["note"] = (
                f"동 전체 평균 경과연수는 {avg_age}년이며, 대상 건물({subject_age}년)은 이보다 "
                f"{cmp_word} 편입니다."
            )
        else:
            result["note"] = f"동 전체 평균 경과연수는 {avg_age}년입니다."
        result["callout"] = {"value": f"{avg_age}년", "label": "동 전체 평균 경과연수"}

    by_purpose = stats.get("주용도별")
    if by_purpose is not None and not by_purpose.empty:
        result["mix"] = [
            {"label": row["주용도"], "value": float(row["비율(%)"])} for _, row in by_purpose.iterrows()
        ]

    benchmark = stats.get("규모벤치마크")
    if benchmark is not None and not benchmark.empty:
        def _fmt(v):
            return "-" if pd.isna(v) else f"{v:g}"
        result["benchmark"] = [
            [row["주용도"], _fmt(row.get("층수(중앙값)")), _fmt(row.get("용적률(중앙값,%)")), _fmt(row.get("높이(중앙값,m)"))]
            for _, row in benchmark.iterrows()
        ]
    return result


def _build_price_history(master, district_price_df=None):
    """공시가격 시계열 + 동단위 주변사례 평균/중앙값/백분위(analyze_district_price_stats).

    동단위 통계는 지번 없이(전체 동) 조회한 district_price_df를 집계만 하므로
    API 재호출 없이 계산되며, 대상 물건이 동 전체에서 상위 몇 %인지까지 붙여
    공시가격 문구를 "주변사례 평균 대비"로 더 정확하게 만든다."""
    ph = master.get("공시가격") or {}
    units = ph.get("단위목록") or []
    if not units:
        return {"trend": [], "rows": [], "note": ph.get("경고"), "district": None}
    unit = units[0]
    timeline = unit.get("추이")
    trend = []
    rows = []
    if timeline is not None and not timeline.empty:
        prev_price = None
        for _, row in timeline.sort_values("연도").iterrows():
            year, price = int(row["연도"]), float(row["주택가격"])
            trend.append({"label": str(year), "value": round(price / 1e8, 2)})
            change = f"{(price / prev_price - 1) * 100:+.1f}%" if prev_price else "-"
            rows.append([str(year), f"{price / 1e8:.1f}억", change])
            prev_price = price
        rows = list(reversed(rows))[:5]
    note = None
    cagr = unit.get("연평균상승률CAGR(%)")
    if cagr is not None:
        note = f"연평균(CAGR) 약 {cagr:.1f}% {'상승' if cagr >= 0 else '하락'}하는 흐름을 보이고 있습니다."

    district = None
    subject_price = unit.get("최신가격")
    subject_price_eok = (subject_price / 1e8) if subject_price else None
    dstats = analyze_district_price_stats(district_price_df, subject_price_eok=subject_price_eok)
    summary = dstats.get("총괄") or {}
    if summary:
        percentile = dstats.get("백분위")
        district = {
            "평균": summary.get("평균공시가격(억)"), "중앙값": summary.get("중앙값공시가격(억)"),
            "호수": summary.get("호수(유닛수)"), "기준연도": summary.get("기준연도"),
            "상위백분율": round(100 - percentile, 1) if percentile is not None else None,
        }
        district_note = (
            f"동 전체(호수 {district['호수']:,}건, {district['기준연도']}년 기준) 주변사례 평균 공시가격은 "
            f"{district['평균']:.2f}억(중앙값 {district['중앙값']:.2f}억)입니다."
        )
        if district["상위백분율"] is not None:
            district_note += f" 대상 물건은 동 전체에서 상위 {district['상위백분율']:.0f}% 수준입니다."
        note = f"{note} {district_note}" if note else district_note

    return {"trend": trend, "rows": rows, "note": note, "district": district}


def _build_commercial(reb_key, sangkwon_key, dong_name, lon, lat, seoul_detail=None):
    """공실률/주변 업종 조회. 실패 사유를 삼키지 않고 result["notes"]에 남겨서
    리포트에 왜 이 섹션이 비었는지 대시보드에서 그대로 보여줄 수 있게 한다.

    업종 Top5·업종 상세표는 "상권분석은 data.seoul.go.kr 중심"이라는 방침에 따라
    seoul_detail(서울 열린데이터광장 우리마을가게 상권분석서비스, _build_seoul()의
    결과)이 있으면 그걸 우선 쓴다. 소상공인시장진흥공단 상가업소 API(전국 대상,
    분류가 거친 업종 소분류 — "펜션" 등 부정확해 보이는 카테고리가 섞여 있음)는
    서울 데이터를 못 쓸 때(비서울 주소, 서울 키 미입력)만 대체용으로 쓴다."""
    result = {
        "vacancy_trend": [], "vacancy_label": "", "top_industries": [], "store_list": [],
        "store_list_cols": None, "area_name": None, "notes": [],
    }

    if seoul_detail and (seoul_detail.get("industry_by_count") or seoul_detail.get("industry_full_table")):
        result["area_name"] = seoul_detail.get("name")
        result["top_industries"] = seoul_detail.get("industry_by_count") or []
        result["store_list"] = seoul_detail.get("industry_full_table") or []
        result["store_list_cols"] = ["업종", "당월매출", "점포수", "개업률", "폐업률"]

    if not reb_key:
        result["notes"].append("공실률: 한국부동산원 인증키 미입력")
    else:
        try:
            statbl_id = REB_COMMERCIAL_VACANCY_STATBL_IDS["중대형 상가"]
            quarter = reb_current_quarter_id()
            snap_df, used_quarter, _ = get_reb_vacancy_snapshot(reb_key, statbl_id, quarter)
            if snap_df is None or snap_df.empty:
                result["notes"].append("공실률: 한국부동산원 응답이 비어 있음")
            else:
                keyword = dong_name[:-1] if dong_name and dong_name.endswith("동") else dong_name
                match = snap_df[snap_df["CLS_FULLNM"].astype(str).str.contains(keyword, na=False)] if keyword else pd.DataFrame()
                if match.empty:
                    result["notes"].append(f"공실률: '{keyword}'이(가) 한국부동산원 상권 분류명과 매칭되지 않음")
                else:
                    cls_id = match.iloc[0]["CLS_ID"]
                    result["vacancy_label"] = f"중대형 상가 · {match.iloc[0]['CLS_FULLNM']}"
                    trend_df, _ = get_reb_vacancy_trend(reb_key, statbl_id, cls_id, "202403", used_quarter)
                    if trend_df is None or trend_df.empty:
                        result["notes"].append("공실률: 추이 데이터 없음")
                    else:
                        trend_df = trend_df.sort_values("WRTTIME_IDTFR_ID")
                        for _, row in trend_df.iterrows():
                            wt = str(row["WRTTIME_IDTFR_ID"])
                            result["vacancy_trend"].append({
                                "label": f"{wt[2:4]}.Q{wt[4:]}",
                                "value": round(float(row["DTA_VAL"]), 1),
                            })
        except Exception as e:
            result["notes"].append(f"공실률: 조회 오류 ({e})")

    if result["top_industries"]:
        pass  # 서울 열린데이터광장 데이터를 이미 채웠으므로 소상공인 API는 건너뜀.
    elif not sangkwon_key:
        result["notes"].append("주변 업종: 소상공인시장진흥공단 인증키 미입력 (서울 열린데이터광장도 사용 불가)")
    elif not (lon and lat):
        result["notes"].append("주변 업종: 좌표(지오코딩) 실패로 조회 불가")
    else:
        try:
            stores_df = get_nearby_stores(sangkwon_key, lon, lat, radius=500)
            if stores_df is None or stores_df.empty or "indsLclsNm" not in stores_df.columns:
                result["notes"].append("주변 업종: 반경 500m 내 데이터 없음")
            else:
                counts = stores_df["indsSclsNm"].value_counts().head(5)
                result["top_industries"] = [{"label": k, "value": int(v)} for k, v in counts.items()]

                top10 = stores_df.head(10)
                result["store_list"] = [
                    [
                        row.get("bizesNm") or "-",
                        row.get("indsSclsNm") or row.get("indsMclsNm") or "-",
                        row.get("rdnmAdr") or row.get("lnoAdr") or "-",
                        f"{int(row['거리(m)'])}m" if pd.notna(row.get("거리(m)")) else "-",
                    ]
                    for _, row in top10.iterrows()
                ]
                result["store_list_cols"] = ["상호명", "업종", "주소", "거리"]
        except Exception as e:
            result["notes"].append(f"주변 업종: 조회 오류 ({e})")

    result["insight"] = None
    vt = result["vacancy_trend"]
    if vt and len(vt) >= 2:
        first, last = vt[0]["value"], vt[-1]["value"]
        trend_word = "안정적으로 유지되고" if abs(last - first) <= 1.5 else ("개선되는" if last < first else "악화되는")
        result["insight"] = f"공실률이 {vt[0]['label']} {first}%에서 {vt[-1]['label']} {last}%로 {trend_word} 흐름입니다."
    elif result["top_industries"]:
        top = result["top_industries"][0]
        result["insight"] = f"업종 구성상 '{top['label']}'이(가) 가장 많은 비중({top['value']}개)을 차지합니다."
    return result


def _build_seoul(seoul_key, lon, lat, dong_name="", locations_loader=None, quarter_loader=None):
    """반환값 3번째는 실패 사유(성공 시 None) — 서울 상권분석/상권영역 지도가 왜 빠졌는지
    대시보드에 그대로 보여주기 위함.

    dong_name(예: "이태원동")을 넘기면 상권명에 그 동 이름이 포함된 상권을
    우선 매칭한다 — 단순 최근접만 쓰면 상권 경계가 촘촘한 지역에서 이름이
    다른 인접 상권이 뽑힐 수 있어, "상권영역이 실제 동과 다르게 나온다"는
    혼선을 줄이기 위함이다.

    locations_loader/quarter_loader를 넘기면(대시보드의 _load_seoul_trade_area_locations/
    _load_seoul_quarter_dataset — @st.cache_data로 6시간 캐시됨) 그걸 대신 쓴다. 사용자가
    "서울 상권분석" 탭을 먼저 조회해뒀다면, 자동 리포트가 같은 데이터를 API 재호출 없이
    캐시로 즉시 재사용하게 하기 위함이다. 안 넘기면 building_example.py의 비캐시 함수로
    대체 동작한다."""
    fetch_locations = locations_loader or get_seoul_trade_area_locations
    fetch_quarter = quarter_loader or get_seoul_trade_area_quarter_dataset
    if not seoul_key:
        return None, None, "서울 열린데이터광장 인증키 미입력"
    if not lon or not lat:
        return None, None, "좌표(지오코딩) 실패로 조회 불가"
    try:
        locations_df = fetch_locations(seoul_key)
        if locations_df is None or locations_df.empty:
            return None, None, "서울시 상권 목록 응답이 비어 있음"
        keyword = dong_name[:-1] if dong_name and dong_name.endswith("동") else dong_name
        trdar_row, dist_m = find_nearest_seoul_trade_area(locations_df, lon, lat, keyword=keyword)
        if dist_m > 1500:
            return None, None, f"가장 가까운 서울시 상권이 {dist_m:.0f}m 떨어져 있어(1.5km 초과) 매칭하지 않음"
        name = trdar_row["TRDAR_CD_NM"]

        selng_df, _ = fetch_quarter(seoul_key, SEOUL_TRDAR_SALES_SERVICE)
        stor_df, _ = fetch_quarter(seoul_key, SEOUL_TRDAR_STORE_SERVICE)
        flpop_df, _ = fetch_quarter(seoul_key, SEOUL_TRDAR_FLPOP_SERVICE)
        wrc_df, _ = fetch_quarter(seoul_key, SEOUL_TRDAR_WRC_POPLTN_SERVICE)

        detail = analyze_seoul_trade_area_detail(trdar_row, selng_df, stor_df, flpop_df, wrc_df)
        s = detail["총괄"]

        stats = []
        if s["추정매출(원)"]:
            stats.append({"label": "추정매출 (당월)", "value": f"{s['추정매출(원)'] / 1e8:,.1f}억원"})
        if s["생활인구"] is not None:
            stats.append({"label": "생활인구 (분기)", "value": f"{s['생활인구']:,.0f}명"})
        if s["직장인구"] is not None:
            stats.append({"label": "직장인구 (분기)", "value": f"{s['직장인구']:,.0f}명"})

        top_industries = []
        industry_by_count = []
        industry_full_table = []
        weekday = []
        industry_df = detail["업종별"]
        if not industry_df.empty:
            top5 = industry_df.head(5)
            top_industries = [
                [
                    r["SVC_INDUTY_CD_NM"], f"{r['THSMON_SELNG_AMT'] / 1e8:.1f}억",
                    f"{int(r['STOR_CO'])}개" if pd.notna(r.get("STOR_CO")) else "-",
                ]
                for _, r in top5.iterrows()
            ]
            if "STOR_CO" in industry_df.columns:
                by_count = industry_df.dropna(subset=["STOR_CO"]).sort_values("STOR_CO", ascending=False).head(5)
                industry_by_count = [
                    {"label": r["SVC_INDUTY_CD_NM"], "value": int(r["STOR_CO"])} for _, r in by_count.iterrows()
                ]
            for _, r in industry_df.head(15).iterrows():
                industry_full_table.append([
                    r["SVC_INDUTY_CD_NM"],
                    f"{r['THSMON_SELNG_AMT'] / 1e8:.1f}억" if pd.notna(r.get("THSMON_SELNG_AMT")) else "-",
                    f"{int(r['STOR_CO'])}개" if pd.notna(r.get("STOR_CO")) else "-",
                    f"{r['OPBIZ_RT']:.1f}%" if pd.notna(r.get("OPBIZ_RT")) else "-",
                    f"{r['CLSBIZ_RT']:.1f}%" if pd.notna(r.get("CLSBIZ_RT")) else "-",
                ])
        if not detail["요일별매출"].empty:
            weekday = [
                {"label": r["요일"], "value": r["매출액(억원)"]} for _, r in detail["요일별매출"].iterrows()
            ]

        gender = detail["성별매출"]
        age_df = detail["연령대별매출"]
        age = [
            {"label": r["연령대"], "value": r["비율(%)"]} for _, r in age_df.iterrows()
        ] if not age_df.empty else []
        weekend_mix = detail["주중주말매출"]
        store = detail["점포"]

        def _population_age_chart(pop_stats):
            age_df = (pop_stats or {}).get("연령대별")
            if age_df is None or age_df.empty:
                return []
            return [{"label": r["연령대"], "value": int(r["인구수"])} for _, r in age_df.iterrows()]

        flpop = detail["생활인구"]
        wrc_pop = detail["직장인구"]

        seoul_detail = {
            "name": name, "stats": stats, "top_industries": top_industries, "weekday": weekday,
            "gender": gender, "age": age, "weekend_mix": weekend_mix, "store": store,
            "insights": detail["인사이트"],
            "industry_by_count": industry_by_count, "industry_full_table": industry_full_table,
            "flpop": flpop, "flpop_age": _population_age_chart(flpop),
            "wrc": wrc_pop, "wrc_age": _population_age_chart(wrc_pop),
        }
        trade_area_map = {"name": name, "lat": trdar_row["lat"], "lon": trdar_row["lon"]}
        return seoul_detail, trade_area_map, None
    except Exception as e:
        return None, None, f"조회 오류 ({e})"


_SEISMIC_SHORT_LABELS = {
    "내진설계 적용(명시)": "적용",
    "내진설계 미적용(명시)": "미적용",
    "준공년도 미상": "미상",
    "1988년 이전 준공 (의무화 전, 미적용 추정)": "미적용(추정)",
    "1988~2017 준공 (기준 완화기간, 확인 필요)": "부분 적용(추정)",
    "2018년 이후 준공 (전면의무화, 적용 추정)": "적용(추정)",
}


def _short_seismic_label(classification: str) -> str:
    return _SEISMIC_SHORT_LABELS.get(classification, str(classification)[:8])


def _build_conclusion(master, data):
    points = []
    old_df = master.get("노후도")
    if old_df is not None and not old_df.empty:
        age = int(old_df.iloc[0]["경과연수"])
        if age >= 30:
            points.append(f"준공 {age}년 경과한 노후 건물로, 매입 시 리모델링/재건축 비용을 사전 검토할 필요가 있습니다.")

    seismic = (master.get("내진분석") or {}).get("취약우선목록")
    if seismic is not None and not seismic.empty and "미적용" in str(seismic.iloc[0].get("내진분류", "")):
        points.append("내진설계가 적용되지 않은 것으로 분류되어, 구조 보강 여부를 확인할 필요가 있습니다.")

    com = data.get("commercial") or {}
    if com.get("vacancy_trend") and len(com["vacancy_trend"]) >= 2:
        first, last = com["vacancy_trend"][0]["value"], com["vacancy_trend"][-1]["value"]
        trend_word = "안정적으로 유지" if abs(last - first) <= 1.5 else ("개선되는" if last < first else "악화되는")
        points.append(f"최근 공실률이 {first}%→{last}%로 {trend_word} 흐름을 보이고 있습니다.")

    ph = data.get("price_history") or {}
    # ph["note"]는 트렌드가 있을 때만 CAGR 요약 문장이고, 데이터가 없으면 일반 면책조항
    # 문구(공시가격≠시세)가 그대로 들어있어 결론 문장으로 쓰기에 부적절하다 — trend 유무로 구분.
    if ph.get("trend") and ph.get("note"):
        points.append(f"공시가격은 {ph['note']}")

    tx = data.get("transactions") or {}
    if tx.get("trend") and ph.get("trend"):
        latest_tx = tx["trend"][-1]["value"]
        latest_ph = ph["trend"][-1]["value"]
        if latest_ph:
            ratio = latest_tx / latest_ph * 100
            points.append(
                f"최근 실거래가({latest_tx:.1f}억)는 최근 공시가격({latest_ph:.1f}억)의 약 {ratio:.0f}% 수준으로, "
                f"{'현실화율 대비 높게' if ratio >= 150 else '통상적인 시세-공시가격 격차 범위 내로'} 형성되어 있습니다."
            )

    if not points:
        points.append("자동 분석 근거가 될 데이터가 충분하지 않아, 추가 리서치가 필요합니다.")
    return points


def fetch_report_data(
    *, service_key, sido, sigungu_name, dong_name, sigungu_code, bdong_code, bun, ji,
    kakao_key=None, vworld_key=None, reb_key=None, sangkwon_key=None, seoul_key=None,
    months_lookback=12, growth_drivers_text=None, swot_text=None, ledger_docs=None,
    progress_callback=None,
    district_title_loader=None, district_price_loader=None,
    seoul_locations_loader=None, seoul_quarter_loader=None,
    single_report=None,
) -> dict:
    """주소(시군구/법정동/번지) 하나에 대해 실제 데이터를 모아 generate_pptx()용 딕셔너리로 반환.

    district_title_loader/district_price_loader/seoul_locations_loader/seoul_quarter_loader는
    대시보드의 @st.cache_data 래퍼(_load_district_titles 등)를 그대로 넘겨받기 위한 훅이다.
    사용자가 동단위통계·노후건축물·공시가격시계열·서울상권분석 탭을 이 리포트보다 먼저
    조회해뒀다면, 같은 캐시를 맞고 API 재호출 없이 즉시 재사용된다 — 안 넘기면(None) 지금처럼
    매번 새로 조회한다.

    single_report는 "종합 리포트" 탭이 이미 조회해 둔 {ledger_type: DataFrame} 결과다.
    호출부(대시보드)가 그 조회에 쓰인 sigungu_code/bdong_code/bun/ji가 지금 이 리포트와
    정확히 같을 때만 넘겨야 한다 — build_master_report()로 그대로 전달돼 표제부·주택가격·
    지역지구구역 재조회를 건너뛴다."""
    from PublicDataReader import BuildingLedger, TransactionPrice

    def _progress(msg):
        if progress_callback:
            progress_callback(msg)

    api = BuildingLedger(service_key)
    tp_api = TransactionPrice(service_key)
    address_label = f"{sido} {sigungu_name} {dong_name}" + (f" {bun}" if bun else "") + (f"-{ji}" if ji and str(ji) != "0" else "")

    _progress("동단위 표제부 조회 중... (탭에서 미리 조회했다면 캐시로 즉시 완료)")
    try:
        if district_title_loader:
            district_title_df = district_title_loader(service_key, sigungu_code, bdong_code)
        else:
            district_title_df = get_building_ledger(
                api, ledger_type="표제부", sigungu_code=sigungu_code, bdong_code=bdong_code,
                max_rows=10000, wait_time=0.15,
            )
    except Exception:
        district_title_df = None

    _progress("동단위 공시가격 조회 중... (탭에서 미리 조회했다면 캐시로 즉시 완료)")
    try:
        if district_price_loader:
            district_price_df = district_price_loader(service_key, sigungu_code, bdong_code)
        else:
            district_price_df = get_building_ledger(
                api, ledger_type="주택가격", sigungu_code=sigungu_code, bdong_code=bdong_code,
                max_rows=10000, wait_time=0.15,
            )
    except Exception:
        district_price_df = None

    _progress("건축물대장·실거래가·공시가격 조회 중...")
    master = build_master_report(
        api, tp_api, sigungu_code, bdong_code, bun, ji,
        months_lookback=months_lookback, district_title_df=district_title_df,
        sido=sido, sigungu_name=sigungu_name, dong_name=dong_name,
        vworld_key=vworld_key, kakao_key=kakao_key, prefetched_ledgers=single_report,
    )

    data = {
        "address": address_label,
        "report_date": datetime.date.today().strftime("%Y년 %m월 %d일 기준"),
    }

    title_df = master.get("표제부")
    core_row = title_df.iloc[0] if title_df is not None and not title_df.empty else None

    # ---- 건축물 개요 ----
    core_fields = [
        ("대지면적", "대지면적", "㎡"), ("건축면적", "건축면적", "㎡"), ("연면적", "연면적", "㎡"),
        ("건폐율", "건폐율", "%"), ("용적률", "용적률", "%"), ("구조코드명", "구조", ""),
        ("지붕코드명", "지붕", ""), ("사용승인일", "사용승인일", ""), ("허가일", "허가일", ""), ("착공일", "착공일", ""),
    ]
    core_list = []
    if core_row is not None:
        for col, label, unit in core_fields:
            val = core_row.get(col)
            if val not in (None, "", "nan") and str(val).strip():
                core_list.append((label, f"{val}{unit}"))
        floors_top = core_row.get("지상층수")
        floors_base = core_row.get("지하층수")
        if floors_top or floors_base:
            core_list.append(("지상/지하층수", f"{floors_top or 0}층 / {floors_base or 0}층"))

    # 위반건축물 여부는 건축HUB Open API에 없는 필드라, 사이드바에서 업로드한 열람본
    # 이미지의 OCR 감지 결과로만 알 수 있다 — 감지되면 개요 카드에 눈에 띄게 추가한다.
    ledger_docs = ledger_docs or []
    if any(doc.get("is_violation") for doc in ledger_docs):
        core_list.append(("위반건축물", "⚠ 있음 (업로드 문서 기준)"))

    zoning = combine_zoning_sources(master)

    _progress("층별 면적 조회 중...")
    floors = []
    try:
        floor_df = get_building_ledger(
            api, ledger_type="층별개요", sigungu_code=sigungu_code, bdong_code=bdong_code,
            bun=bun, ji=ji, max_rows=200, wait_time=0.15,
        )
        if floor_df is not None and not floor_df.empty and "면적" in floor_df.columns and "층번호명" in floor_df.columns:
            floor_df = floor_df.copy()
            floor_df["_면적"] = pd.to_numeric(floor_df["면적"], errors="coerce")
            floor_df = floor_df[floor_df["_면적"].notna() & (floor_df["_면적"] > 0)]
            floors = [
                {"label": row["층번호명"], "value": round(float(row["_면적"]), 1)}
                for _, row in floor_df.head(20).iterrows()
            ]
    except Exception:
        pass
    data["building"] = {"core": core_list, "floors": floors, "zoning": zoning}

    # ---- 위치 및 입지 ----
    coord = master.get("좌표")
    lon, lat = (coord if coord else (None, None))
    subway = _nearby_subway(lat, lon) if lat and lon else []
    map_buf = _render_location_map(lat, lon, subway) if lat and lon else None
    data["location"] = {
        "adong_name": dong_name, "ldong_name": dong_name,
        "subway": subway, "map_image": map_buf,
    }
    data["lon"], data["lat"] = lon, lat

    # ---- 핵심 요약 ----
    data["summary_text"] = build_executive_summary(master)
    summary_stats = []
    old_df = master.get("노후도")
    if core_row is not None:
        approval_year = extract_year(core_row.get("사용승인일"))
        age = int(old_df.iloc[0]["경과연수"]) if old_df is not None and not old_df.empty else None
        if approval_year:
            sub = f"{age}년 경과" + (" · 노후" if age and age >= 30 else "") if age is not None else None
            summary_stats.append({"label": "사용승인일 (경과연수)", "value": f"{approval_year}년", "sub": sub})
    seismic_list = (master.get("내진분석") or {}).get("취약우선목록")
    seismic_data = {"dong_dist": [], "age_dist": [], "subject_label": None, "subject_full": None}
    if seismic_list is not None and not seismic_list.empty:
        seismic_full = str(seismic_list.iloc[0]["내진분류"])
        seismic_data["subject_label"] = _short_seismic_label(seismic_full)
        seismic_data["subject_full"] = seismic_full
        summary_stats.append({"label": "내진 설계", "value": _short_seismic_label(seismic_full), "sub": seismic_full})

    # 대상 건물 내진분석(master["내진분석"])은 지번 하나(호실 1건)짜리라 분포를 그릴 수 없다 —
    # 동 전체 표제부(district_title_df, 위에서 이미 조회함)를 따로 analyze_seismic_risk에
    # 태워서 동단위 분류 분포를 얻는다. API 재호출 없이 로컬 재계산만 한다.
    try:
        if district_title_df is not None and not district_title_df.empty:
            dong_seismic = analyze_seismic_risk(district_title_df, top_n=1)
            dist_df = dong_seismic.get("분류별집계")
            if dist_df is not None and not dist_df.empty:
                seismic_data["dong_dist"] = [
                    {"label": _short_seismic_label(row["분류"]), "value": int(row["건수"])}
                    for _, row in dist_df.iterrows()
                ]
    except Exception:
        pass
    age_dist_df = (master.get("동단위통계") or {}).get("노후도분포")
    if age_dist_df is not None and not age_dist_df.empty:
        seismic_data["age_dist"] = [
            {"label": row["구간"], "value": int(row["건수"])} for _, row in age_dist_df.iterrows()
        ]
    data["seismic"] = seismic_data

    # 어떤 섹션이 왜 빠졌는지(키 미입력/데이터 없음/API 오류) 대시보드에 그대로 보여주기 위한 메모.
    notes = []
    if not lon or not lat:
        notes.append("위치·좌표: 지오코딩 실패 — 카카오맵/브이월드 키를 확인하세요 (지도·주변 상가업소·서울 상권분석에 영향)")

    _progress("동단위 시장 통계 분석 중...")
    data["district"] = _build_district(master)
    if not data["district"]["age_buckets"]:
        notes.append("동단위 시장 통계: 동 표제부 데이터를 가져오지 못했거나 비어 있음")

    _progress("실거래가 정리 중...")
    data["transactions"] = _build_transactions(master, months_lookback)
    if data["transactions"]["rows"]:
        latest = data["transactions"]["rows"][0]
        summary_stats.append({"label": "최근 실거래가", "value": latest[3], "sub": f"{latest[1]} · {latest[2]} · {latest[0]}"})
    else:
        notes.append(f"실거래가 동향: 최근 {months_lookback}개월 내 거래 내역 없음")

    _progress("공시가격 시계열 · 동단위 주변사례 분석 중...")
    data["price_history"] = _build_price_history(master, district_price_df)
    if not data["price_history"]["trend"]:
        notes.append("공시가격 시계열: 시계열 데이터 없음")

    # 서울 상권분석을 상권 개황(commercial)보다 먼저 구해서, 업종 Top5/상세표를
    # (전국 대상이라 분류가 거친 소상공인 API 대신) 서울 열린데이터광장 데이터로
    # 채울 수 있게 한다 — "상권분석은 data.seoul.go.kr 중심" 방침.
    is_seoul = sido.startswith("서울")
    seoul_detail, trade_area_map = None, None
    if not is_seoul:
        notes.append("서울 상권분석·상권영역 지도: 서울 열린데이터광장 API는 서울 소재 주소만 지원 (해당 없음)")
    elif not seoul_key:
        notes.append("서울 상권분석·상권영역 지도: 서울 열린데이터광장 인증키 미입력")
    else:
        _progress("서울 상권분석 데이터 조회 중... (탭에서 미리 조회했다면 캐시로 즉시 완료, 아니면 최대 1분 정도)")
        seoul_detail, trade_area_map, seoul_reason = _build_seoul(
            seoul_key, lon, lat, dong_name=dong_name,
            locations_loader=seoul_locations_loader, quarter_loader=seoul_quarter_loader,
        )
        if seoul_reason:
            notes.append(f"서울 상권분석·상권영역 지도: {seoul_reason}")
    data["seoul_trade_area"] = seoul_detail
    if seoul_detail and seoul_detail.get("stats"):
        summary_stats.append({"label": "상권 추정매출(월)", "value": seoul_detail["stats"][0]["value"], "sub": seoul_detail["name"]})

    if trade_area_map:
        map_buf2 = _render_location_map(trade_area_map["lat"], trade_area_map["lon"])
        data["trade_area_map"] = {"name": trade_area_map["name"], "map_image": map_buf2}
    else:
        data["trade_area_map"] = None

    _progress("상업용부동산 공실률 · 주변 업종 조회 중...")
    data["commercial"] = _build_commercial(reb_key, sangkwon_key, dong_name, lon, lat, seoul_detail=seoul_detail)
    notes.extend(data["commercial"].get("notes") or [])
    if data["commercial"]["vacancy_trend"]:
        summary_stats.append({"label": "공실률", "value": f"{data['commercial']['vacancy_trend'][-1]['value']}%", "sub": data["commercial"]["vacancy_label"]})

    data["summary_stats"] = summary_stats

    # 표지에는 상황에 따라 있을 수도 없을 수도 있는 5개 항목 중 "가장 눈에 띄는" 3개만
    # 우선순위대로 골라 쓴다 (예전에는 summary_stats[:3]으로 그냥 앞 3개를 썼는데, 계산
    # 순서에 따라 달라지는 우연의 결과였다).
    _cover_priority = ["최근 실거래가", "내진 설계", "공실률", "상권 추정매출(월)", "사용승인일 (경과연수)"]
    _stats_by_label = {s["label"]: s for s in summary_stats}
    cover_stats = [_stats_by_label[label] for label in _cover_priority if label in _stats_by_label][:3]
    data["cover_stats"] = cover_stats or summary_stats[:3]

    data["notes"] = notes

    def _parse_lines(text):
        return [line.strip() for line in (text or "").splitlines() if line.strip()]

    data["growth_drivers"] = _parse_lines(growth_drivers_text)
    swot_text = swot_text or {}
    data["swot"] = {
        "strengths": _parse_lines(swot_text.get("strengths")),
        "weaknesses": _parse_lines(swot_text.get("weaknesses")),
        "opportunities": _parse_lines(swot_text.get("opportunities")),
        "threats": _parse_lines(swot_text.get("threats")),
    }

    _progress("종합 의견 정리 중...")
    data["conclusion"] = _build_conclusion(master, data)

    data["ledger_content"] = extract_ledger_content(ledger_docs) if ledger_docs else None
    if ledger_docs:
        data["notes"].append(
            f"업로드한 건축물대장 열람본 {len(ledger_docs)}장에서 OCR로 소유자현황·변동사항·위반건축물 여부를 "
            "추출해 리포트 끝의 '소유자현황 · 변동사항' 슬라이드로 정리했습니다. 표 셀 인식 특성상 오탈자가 "
            "있을 수 있어 정확한 값은 대시보드에 표시된 원본 이미지와 대조하는 것을 권장합니다."
        )

    return data


# ==================================================================
# 종합 리포트 탭 전용 pptx — 11종 건축물대장(report dict)을 자동pptx 리포트와
# 같은 디자인 시스템으로 여러 장짜리 문서로 그린다. 예전엔 fpdf2로 항목을
# 나열한 1페이지짜리 흑백 표(generate_pdf_report, building_example.py)였는데,
# "자동pptx 리포트 탭처럼 전문적으로" 요청받아 같은 네이비/테라코타 톤·카드·
# 차트 컴포넌트를 재사용해 별도 생성기로 새로 만들었다. fetch_report_data()가
# 쓰는 market-report용 data dict/_slide_plan 구조와는 무관한, report(=ledger_type
# -> DataFrame) 하나만 입력으로 받는 독립된 흐름이다.
# ==================================================================
def _ledger_core_row(report):
    for key in ("표제부", "총괄표제부", "기본개요"):
        df = report.get(key)
        if df is not None and not df.empty and "오류" not in df.columns:
            return df.iloc[0]
    return None


def _ledger_cover_slide(prs, address_label, core_row, report_date, has_violation=False, violation_detail=None):
    slide = _new_slide(prs, dark=True)
    for cx, cy, cw, ch, color, grad in [
        (8.6, 3.0, 7.5, 7.5, NAVY_LIGHT, False),
        (9.6, 4.0, 5.5, 5.5, NAVY_DARK, False),
        (10.4, 4.8, 3.9, 3.9, None, True),
    ]:
        ellipse = slide.shapes.add_shape(MSO_SHAPE.OVAL, _sx(cx), _sy(cy), _su(cw), _su(ch))
        if grad:
            _gradient_fill(ellipse, angle=45)
        else:
            ellipse.fill.solid()
            ellipse.fill.fore_color.rgb = color
        ellipse.line.fill.background()
        ellipse.shadow.inherit = False

    _textbox(slide, 0.9, 2.35, 8, 0.4, "건축물대장 종합 리포트", size=14, bold=True, color=TERRACOTTA)
    _textbox(slide, 0.9, 2.8, 9.5, 1.6, address_label, size=36, bold=True, color=WHITE)
    _textbox(slide, 0.9, 4.05, 9, 0.5, "표제부 · 층별개요 · 지역지구구역 등 공부 원본 데이터 기준", size=16, color=ICE)
    if has_violation:
        _textbox(slide, 0.9, 4.55, 9, 0.4, "🚩 업로드한 열람본에서 위반건축물 표기 감지됨", size=13, bold=True, color=TERRACOTTA)
        if violation_detail:
            _textbox(slide, 0.9, 4.95, 9.2, 0.5, f"내용(OCR): {violation_detail}", size=11, color=ICE, line_spacing=1.3)

    stats = []
    if core_row is not None:
        if _clean(core_row.get("연면적")):
            stats.append((f"{_clean(core_row.get('연면적'))}㎡", "연면적"))
        if _clean(core_row.get("사용승인일")):
            stats.append((_clean(core_row.get("사용승인일")), "사용승인일"))
        if _clean(core_row.get("용적률")):
            stats.append((f"{_clean(core_row.get('용적률'))}%", "용적률"))
    for i, (value, label) in enumerate(stats[:3]):
        cx = 0.9 + i * 2.9
        _textbox(slide, cx, 5.5, 2.6, 0.55, value, size=22, bold=True, color=TERRACOTTA)
        _textbox(slide, cx, 6.05, 2.6, 0.5, label, size=10.5, color=ICE)

    _textbox(slide, 0.9, 6.7, 6, 0.4, report_date, size=11, color=MUTED)
    return slide


def _ledger_overview_slide(prs, address_label, core_row, report, has_violation=False):
    slide = _new_slide(prs)
    _section_title(slide, "건축물 개요")

    left_x, left_y, left_w, left_h = 0.6, 1.3, 6.7, 5.6
    _card(slide, left_x, left_y, left_w, left_h)
    core = []
    if core_row is not None:
        for col, label in _CORE_FIELD_LABELS:
            val = _clean(core_row.get(col))
            if val:
                core.append((label, val))
    if has_violation:
        core.append(("위반건축물", "⚠ 있음 (업로드 문서 기준)"))
    if core:
        row_h = (left_h - 0.4) / max(1, math.ceil(len(core) / 2))
        for i, (label, value) in enumerate(core):
            col, row = i % 2, i // 2
            x = left_x + 0.3 + col * (left_w / 2 - 0.15)
            y = left_y + 0.25 + row * row_h
            _textbox(slide, x, y, left_w / 2 - 0.5, row_h * 0.42, label, size=10.5, color=MUTED)
            _textbox(slide, x, y + row_h * 0.38, left_w / 2 - 0.5, row_h * 0.5, value, size=14, bold=True, color=NAVY)
    else:
        _textbox(slide, left_x + 0.3, left_y + 0.3, left_w - 0.6, 0.6, "표제부 데이터가 없습니다.", size=12, color=MUTED)

    right_x, right_w = 7.5, 5.2
    next_y = 1.3
    zoning_df = report.get("지역지구구역")
    if zoning_df is not None and not zoning_df.empty and "오류" not in zoning_df.columns and "지역지구구역코드명" in zoning_df.columns:
        zoning = [z for z in (_clean(v) for v in zoning_df["지역지구구역코드명"]) if z]
        if zoning:
            _textbox(slide, right_x, next_y, right_w, 0.35, "용도지역 · 지구", size=13, bold=True, color=NAVY)
            py = next_y + 0.4
            for z in zoning[:5]:
                pw = min(right_w, 0.35 + len(z) * 0.16)
                pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(right_x), _sy(py), _sx(pw), _sy(0.42))
                pill.adjustments[0] = 0.5
                pill.fill.solid()
                pill.fill.fore_color.rgb = ICE
                pill.line.fill.background()
                pill.shadow.inherit = False
                tf = pill.text_frame
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = z
                run.font.size = Pt(10.5)
                run.font.name = FONT_NAME
                run.font.color.rgb = NAVY
                py += 0.55
            next_y = py + 0.3

    one_liners = []
    for ledger_type, cols in _ONE_LINE_FIELDS.items():
        df = report.get(ledger_type)
        if df is None or df.empty or "오류" in df.columns:
            continue
        row = df.iloc[0]
        parts = [f"{c}: {_clean(row[c])}" for c in cols if c in df.columns and _clean(row[c])]
        if parts:
            one_liners.append(f"{ledger_type} — " + " · ".join(parts))
    if one_liners:
        _textbox(slide, right_x, next_y, right_w, 0.35, "기타 공부 정보", size=13, bold=True, color=NAVY)
        _textbox(slide, right_x, next_y + 0.4, right_w, 2.0, "\n".join(one_liners), size=11, line_spacing=1.35)

    _page_footer(slide, address_label, "건축물 개요")


def _ledger_floor_slide(prs, address_label, report):
    df = report.get("층별개요")
    if df is None or df.empty or "오류" in df.columns or "면적" not in df.columns:
        return False
    d = df.copy()
    d["_면적"] = pd.to_numeric(d["면적"], errors="coerce")
    d = d[d["_면적"].notna() & (d["_면적"] > 0)]
    if d.empty:
        return False

    slide = _new_slide(prs)
    _section_title(slide, "층별현황")
    cats = [str(v) for v in d["층번호명"]] if "층번호명" in d.columns else [str(i + 1) for i in range(len(d))]
    vals = [round(float(v), 1) for v in d["_면적"]]
    chart = _bar_chart(slide, 0.6, 1.3, 6.2, 5.6, cats, vals, TERRACOTTA, horizontal=True)
    _gradient_fill(chart.plots[0].series[0].format, angle=0)

    cols = [c for c in _MULTIROW_KEEP_COLS["층별개요"] if c in d.columns]
    rows = [[_clean(v) for v in row] for row in d[cols].itertuples(index=False)]
    _table(slide, 7.1, 1.3, 5.6, min(5.6, 0.5 * (len(rows) + 1)), cols, rows, font_size=10.5)
    _page_footer(slide, address_label, "층별현황")
    return True


def _ledger_multirow_slide(prs, address_label, title, df, cols):
    cols = cols or list(df.columns)
    rows = [[_clean(v) for v in row] for row in df[cols].itertuples(index=False)]
    slide = _new_slide(prs)
    _section_title(slide, title)
    _table(slide, 0.6, 1.3, 12.1, min(5.6, 0.5 * (len(rows) + 1)), cols, rows, font_size=11.5)
    _page_footer(slide, address_label, title)


def _ledger_owner_history_slide(prs, address_label, owners, changes, firms=None, builder=None,
                                 violation_detail=None):
    """업로드한 열람본 이미지에서 OCR로 뽑은 소유자현황·변동사항·관계업체를, 원본 사진
    첨부 대신 이 리포트의 다른 슬라이드와 같은 표/텍스트 형식으로 정리해서 보여준다.
    표 셀 인식이 무너진 OCR 결과라 이름 등이 오탈자로 나올 수 있어(예: '윤명분'->'륜령분')
    그 사실을 슬라이드에 직접 명시한다 — 원본 이미지는 대시보드 화면 쪽에서 항상 대조 가능하다."""
    slide = _new_slide(prs)
    _section_title(slide, "소유자현황 · 변동사항")
    _textbox(
        slide, 0.6, 1.05, 11.5, 0.4,
        "업로드한 건축물대장 열람본에서 OCR로 추출 — 공공 API에는 없는 정보이며, 표 인식 특성상 "
        "오탈자가 있을 수 있어 정확한 값은 원본 이미지 대조를 권장합니다.",
        size=10.5, color=MUTED, line_spacing=1.3,
    )

    left_x, left_w = 0.6, 5.6
    _textbox(slide, left_x, 1.6, left_w, 0.35, "소유자현황", size=13, bold=True, color=NAVY)
    if owners:
        rows = [[o["성명"], o["변동일"]] for o in owners]
        owner_table_h = min(3.0, 0.5 * (len(rows) + 1))
        _table(slide, left_x, 2.0, left_w, owner_table_h, ["성명", "변동일"], rows,
               col_ratios=[2, 1], font_size=12)
        next_left_y = 2.0 + owner_table_h + 0.35
    else:
        _card(slide, left_x, 2.0, left_w, 1.0)
        _textbox(slide, left_x + 0.25, 2.2, left_w - 0.5, 0.6,
                 "OCR로 소유자 정보를 특정하지 못했습니다. 원본 이미지를 확인하세요.", size=11, color=MUTED)
        next_left_y = 3.3

    firms = firms or []
    if firms or builder:
        _textbox(slide, left_x, next_left_y, left_w, 0.35, "설계 · 감리 등 관계자", size=13, bold=True, color=NAVY)
        lines = ([f"건축주: {builder}"] if builder else []) + [f"· {f}" for f in firms]
        card_h = min(2.0, 0.4 + 0.4 * len(lines))
        _card(slide, left_x, next_left_y + 0.4, left_w, card_h)
        _textbox(slide, left_x + 0.25, next_left_y + 0.55, left_w - 0.5, card_h - 0.3,
                 "\n".join(lines), size=11, line_spacing=1.4)

    right_x, right_w = 6.5, 6.2
    right_y = 1.6
    if violation_detail:
        _textbox(slide, right_x, right_y, right_w, 0.3, "🚩 위반건축물 내용", size=13, bold=True, color=TERRACOTTA)
        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _sx(right_x), _sy(right_y + 0.35),
                                          _sx(right_w), _sy(0.9))
        callout.adjustments[0] = 0.08
        callout.fill.solid()
        callout.fill.fore_color.rgb = RGBColor(0xFB, 0xEA, 0xE3)
        callout.line.color.rgb = TERRACOTTA
        callout.line.width = Pt(1)
        callout.shadow.inherit = False
        tf = callout.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.15)
        tf.margin_top = tf.margin_bottom = Inches(0.08)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = violation_detail
        run.font.size = Pt(10.5)
        run.font.name = FONT_NAME
        run.font.color.rgb = TEXT_DARK
        right_y += 0.35 + 0.9 + 0.3

    _textbox(slide, right_x, right_y, right_w, 0.35, "변동사항", size=13, bold=True, color=NAVY)
    if changes:
        card_h = min(7.0 - right_y - 0.4, 0.5 + 0.55 * len(changes))
        _card(slide, right_x, right_y + 0.4, right_w, card_h)
        _textbox(slide, right_x + 0.25, right_y + 0.55, right_w - 0.5, card_h - 0.3,
                 "\n".join(f"· {c}" for c in changes), size=10, line_spacing=1.4)
    else:
        _card(slide, right_x, right_y + 0.4, right_w, 1.0)
        _textbox(slide, right_x + 0.25, right_y + 0.6, right_w - 0.5, 0.6,
                 "OCR로 변동사항을 특정하지 못했습니다. 원본 이미지를 확인하세요.", size=11, color=MUTED)

    _page_footer(slide, address_label, "소유자현황 · 변동사항")


_LEDGER_DATA_SOURCES = [
    "국토교통부 건축HUB 건축물대장정보 서비스 (표제부 · 총괄표제부 · 층별개요 · 지역지구구역 · "
    "오수정화시설 · 주택가격 · 부속지번 · 전유공용면적 등)",
]

_LEDGER_APPENDIX_DISCLAIMER = (
    "· 본 리포트는 건축HUB Open API 응답을 자동 정리한 참고 자료이며, 법적 효력이 있는 건축물대장 발급 문서가 아닙니다.\n"
    "· 위반건축물 여부 · 소유자현황 · 변동사항 · 설계/감리 등 관계자는 이 공공 API에 해당 필드가 없어 표시되지 않습니다"
    "(소유자 오픈API는 2026년 기준 data.go.kr에서 서비스 종료). 정확한 확인은 정부24 건축물대장 열람/발급을 이용하세요.\n"
    "· 사이드바에 열람본 이미지를 업로드했다면, 위 항목은 이 리포트의 '소유자현황 · 변동사항' 슬라이드에서 OCR로 "
    "정리한 값을 볼 수 있습니다 — 표 셀 인식 특성상 오탈자가 있을 수 있으니, 정확한 값은 대시보드에 표시된 원본 "
    "이미지와 대조하세요."
)


def _ledger_appendix_slide(prs, address_label):
    slide = _new_slide(prs)
    _section_title(slide, "부록 — 데이터 출처 및 유의사항")
    _textbox(slide, 0.6, 1.3, 11.5, 0.35, "데이터 출처", size=13, bold=True, color=NAVY)
    _textbox(slide, 0.6, 1.7, 11.5, 1.0, "\n".join(f"· {s}" for s in _LEDGER_DATA_SOURCES), size=12, line_spacing=1.35)
    _textbox(slide, 0.6, 3.0, 11.5, 0.35, "유의사항", size=13, bold=True, color=NAVY)
    _textbox(slide, 0.6, 3.4, 11.5, 2.6, _LEDGER_APPENDIX_DISCLAIMER, size=11, color=MUTED, line_spacing=1.4)
    _page_footer(slide, address_label, "부록")


def generate_ledger_pptx(report: dict, address_label: str, ledger_docs=None) -> bytes:
    """종합 리포트 탭(11종 건축물대장 조회 결과)을 여러 장짜리 pptx로 그려 바이트로 반환."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    core_row = _ledger_core_row(report)
    report_date = datetime.date.today().strftime("%Y년 %m월 %d일 기준")
    content = extract_ledger_content(ledger_docs) if ledger_docs else None
    has_violation = content["is_violation"] if content else False
    violation_detail = content.get("violation_detail") if content else None

    _ledger_cover_slide(prs, address_label, core_row, report_date,
                        has_violation=has_violation, violation_detail=violation_detail)
    _ledger_overview_slide(prs, address_label, core_row, report, has_violation=has_violation)
    _ledger_floor_slide(prs, address_label, report)

    for ledger_type in ("전유공용면적", "부속지번"):
        df = report.get(ledger_type)
        if df is not None and not df.empty and "오류" not in df.columns:
            cols = [c for c in _MULTIROW_KEEP_COLS.get(ledger_type, []) if c in df.columns] or list(df.columns)
            _ledger_multirow_slide(prs, address_label, ledger_type, df, cols)

    if content:
        _ledger_owner_history_slide(
            prs, address_label, content["owners"], content["changes"],
            firms=content.get("firms"), builder=content.get("builder"),
            violation_detail=violation_detail,
        )

    _ledger_appendix_slide(prs, address_label)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
